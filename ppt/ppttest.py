from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random
from enum import Enum, auto

class TextLevel(Enum):
    """文本层级枚举"""
    TITLE = auto()      # 主标题
    SUBTITLE = auto()   # 副标题
    HEADING1 = auto()   # 一级标题
    HEADING2 = auto()   # 二级标题
    HEADING3 = auto()   # 三级标题
    BODY = auto()       # 正文
    CAPTION = auto()    # 说明文字

# 默认字号配置
DEFAULT_FONT_SIZES = {
    TextLevel.TITLE: Pt(36),
    TextLevel.SUBTITLE: Pt(24),
    TextLevel.HEADING1: Pt(28),
    TextLevel.HEADING2: Pt(22),
    TextLevel.HEADING3: Pt(18),
    TextLevel.BODY: Pt(14),
    TextLevel.CAPTION: Pt(12)
}

class FontDesign:
    """增强版字体设计模块"""
    
    # 字体层级配置
    LEVEL_STYLES = {
        TextLevel.TITLE: {
            "font_name": "微软雅黑",
            "color": (0, 0, 0),
            "bold": True,
            "align": PP_ALIGN.CENTER
        },
        TextLevel.SUBTITLE: {
            "font_name": "微软雅黑",
            "color": (100, 100, 100),
            "bold": False,
            "align": PP_ALIGN.CENTER
        },
        TextLevel.HEADING1: {
            "font_name": "微软雅黑",
            "color": (0, 0, 139),  # 深蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.HEADING2: {
            "font_name": "微软雅黑",
            "color": (0, 70, 130),  # 中蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.HEADING3: {
            "font_name": "微软雅黑",
            "color": (0, 100, 120),  # 淡蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.BODY: {
            "font_name": "宋体",
            "color": (0, 0, 0),
            "bold": False,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.CAPTION: {
            "font_name": "宋体",
            "color": (100, 100, 100),
            "bold": False,
            "align": PP_ALIGN.CENTER
        }
    }
    
    @staticmethod
    def set_font(text_frame, level=None, font_name=None, size=None, color=None, bold=None, align=None):
        """
        设置字体样式
        :param level: 文本层级
        :param font_name, size, color, bold, align: 可选参数，覆盖默认设置
        """
        paragraph = text_frame.paragraphs[0]
        font = paragraph.font
        
        # 如果指定了level，使用预定义样式
        if level is not None:
            style = FontDesign.LEVEL_STYLES.get(level, {})
            font.name = font_name or style.get("font_name", "微软雅黑")
            font.size = size or DEFAULT_FONT_SIZES.get(level, Pt(12))
            font.bold = bold if bold is not None else style.get("bold", False)
            font.color.rgb = RGBColor(*(color or style.get("color", (0, 0, 0))))
            paragraph.alignment = align or style.get("align", PP_ALIGN.LEFT)
        else:
            # 兼容旧接口
            font.name = font_name or "微软雅黑"
            font.size = size or Pt(12)
            font.bold = bold if bold is not None else False
            font.color.rgb = RGBColor(*(color or (0, 0, 0)))
            paragraph.alignment = align or PP_ALIGN.LEFT

class TextBoxDesign:
    """智能文本框设计模块"""
    
    # 静态类属性，存储当前幻灯片的尺寸
    current_slide_width = 13.33  # 默认16:9宽度（英寸）
    current_slide_height = 7.5   # 默认16:9高度（英寸）
    
    # 预定义的位置配置（以幻灯片尺寸的比例表示）
    POSITION_PROFILES = {
        "header": (0.2, 0.1, 0.6, 0.15),        # 页眉区域
        "left_sidebar": (0.05, 0.3, 0.25, 0.6), # 左侧边栏
        "right_sidebar": (0.7, 0.3, 0.25, 0.6), # 右侧边栏
        "main_content": (0.3, 0.3, 0.65, 0.6),  # 主内容区
        "footer": (0.2, 0.85, 0.6, 0.1)         # 页脚区域
    }
    
    @staticmethod
    def _estimate_lines(text, font_size, box_width):
        """估算文本行数（按字符数简单计算）"""
        # 防止除以零错误
        box_width_pt = box_width * 72  # 将英寸转换为点
        chars_per_pt = 0.6  # 经验值：0.6pt/字符
        divisor = font_size.pt * chars_per_pt
        if divisor <= 0:
            avg_chars_per_line = 30  # 默认值，如果计算结果无效
        else:
            avg_chars_per_line = max(1, int(box_width_pt / divisor))
        return max(1, len(text) // avg_chars_per_line)

    @staticmethod
    def add_textbox(slide, left=None, top=None, width=None, height=None, 
                    text="", level=None, position_profile=None, 
                    auto_adjust=True, min_size_ratio=0.6, **font_kwargs):
        """
        添加智能文本框
        :param slide: 幻灯片对象
        :param left, top, width, height: 位置和尺寸（英寸）
        :param text: 文本内容
        :param level: 文本层级
        :param position_profile: 位置配置名称
        :param auto_adjust: 是否自动调整字号
        :param min_size_ratio: 最小可接受字号比例(相对于默认字号)
        :param font_kwargs: 其他字体参数
        """
        # 处理位置配置
        if position_profile:
            try:
                # 使用类属性存储的幻灯片尺寸
                slide_width = TextBoxDesign.current_slide_width
                slide_height = TextBoxDesign.current_slide_height
                
                # 计算实际位置和尺寸
                profile = TextBoxDesign.POSITION_PROFILES.get(position_profile)
                if profile:
                    left_ratio, top_ratio, width_ratio, height_ratio = profile
                    left = left or (slide_width * left_ratio)
                    top = top or (slide_height * top_ratio)
                    width = width or (slide_width * width_ratio)
                    height = height or (slide_height * height_ratio)
            except Exception as e:
                print(f"警告：无法应用位置配置'{position_profile}'，使用默认位置。错误：{str(e)}")
        
        # 确保所有位置参数存在
        if None in (left, top, width, height):
            # 默认值
            left = left or 1
            top = top or 1
            width = width or 8
            height = height or 1
        
        # 创建文本框 - 使用全局导入的Inches
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.text = text
        
        # 获取默认字号
        default_size = DEFAULT_FONT_SIZES.get(level, Pt(12)) if level else font_kwargs.get('size', Pt(12))
        
        # 设置默认字体
        if level:
            FontDesign.set_font(text_frame, level=level, **font_kwargs)
        else:
            FontDesign.set_font(text_frame, **font_kwargs)
        
        # 自动调整字号
        if auto_adjust and text:
            # 计算当前字号下的行数
            current_size = default_size
            font_size = default_size
            min_size = Pt(max(8, default_size.pt * min_size_ratio))  # 最小不低于8pt
            
            estimated_lines = TextBoxDesign._estimate_lines(text, font_size, width)
            max_lines = int(height * 72 / (font_size.pt * 1.2))  # 1.2为行距系数
            
            # 如果超幅，减小字号
            while estimated_lines > max_lines and font_size.pt > min_size.pt:
                font_size = Pt(font_size.pt - 1)
                estimated_lines = TextBoxDesign._estimate_lines(text, font_size, width)
                max_lines = int(height * 72 / (font_size.pt * 1.2))
            
            # 如果字号有调整，更新字体
            if font_size.pt != current_size.pt:
                print(f"信息：文本框（位置：{left},{top}）字号从{current_size.pt}pt自动调整为{font_size.pt}pt")
                if level:
                    FontDesign.set_font(text_frame, level=level, size=font_size, **font_kwargs)
                else:
                    FontDesign.set_font(text_frame, size=font_size, **font_kwargs)
            
            # 如果仍然超幅，发出警告
            if estimated_lines > max_lines:
                print(f"警告：文本框可能超幅（位置：{left},{top}，预估行数：{estimated_lines} > 最大行数：{max_lines}）")
        
        return textbox

    @staticmethod
    def add_paragraph(text_frame, text, level=None, **font_kwargs):
        """
        添加新段落并保持样式一致
        :param text_frame: 文本框架
        :param text: 段落文本
        :param level: 文本层级
        :param font_kwargs: 字体参数
        """
        paragraph = text_frame.add_paragraph()
        paragraph.text = text
        
        if level:
            font_style = FontDesign.LEVEL_STYLES.get(level, {})
            for idx, run in enumerate(paragraph.runs):
                run.font.name = font_kwargs.get('font_name', font_style.get('font_name', '微软雅黑'))
                run.font.size = font_kwargs.get('size', DEFAULT_FONT_SIZES.get(level, Pt(12)))
                run.font.bold = font_kwargs.get('bold', font_style.get('bold', False))
                run.font.color.rgb = RGBColor(*font_kwargs.get('color', font_style.get('color', (0, 0, 0))))
            paragraph.alignment = font_kwargs.get('align', font_style.get('align', PP_ALIGN.LEFT))
        else:
            # 使用传入的参数
            for idx, run in enumerate(paragraph.runs):
                run.font.name = font_kwargs.get('font_name', '微软雅黑')
                run.font.size = font_kwargs.get('size', Pt(12))
                run.font.bold = font_kwargs.get('bold', False)
                run.font.color.rgb = RGBColor(*font_kwargs.get('color', (0, 0, 0)))
            paragraph.alignment = font_kwargs.get('align', PP_ALIGN.LEFT)
        
        return paragraph

class SlideHeaderFooter:
    """页眉页脚设计模块"""
    @staticmethod
    def set_header_footer(slide, header_text="", footer_text="", footer_font_size=Pt(10)):
        # 页眉（需通过修改XML实现）
        if header_text:
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 12:  # 12对应页眉
                    shape.text = header_text

        # 页脚（简单实现）
        if footer_text:
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8.0), Inches(0.5))
            footer.text = footer_text
            FontDesign.set_font(footer.text_frame, size=footer_font_size, align=PP_ALIGN.CENTER)

class SlideSizeDesign:
    """幻灯片尺寸设计模块"""
    @staticmethod
    def set_slide_size(presentation, width=Inches(13.33), height=Inches(7.5)):
        presentation.slide_width = width
        presentation.slide_height = height

class ArtTextDesign:
    """艺术字设计模块（改进版）"""
    @staticmethod
    def add_art_text(
        slide, left, top, width, height, text="ART", 
        font_size=Pt(48), use_shape=False, 
        shape_color=(240, 240, 240), text_color=None,
        **font_kwargs
    ):
        """
        添加艺术字
        :param use_shape: 是否使用形状背景
        :param shape_color: 背景色RGB值
        :param text_color: 文字颜色RGB值（None则随机）
        """
        if use_shape:
            # 带背景框模式
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.text = text
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*shape_color)
            shape.line.color.rgb = RGBColor(150, 150, 150)
        else:
            # 纯文字模式（普通文本框）
            shape = TextBoxDesign.add_textbox(
                slide, left, top, width, height, text,
                check_overflow=False,  # 艺术字通常不需要检测超幅
                **font_kwargs
            )
        
        # 统一设置文字样式
        text_frame = shape.text_frame
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        color = text_color or (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
        FontDesign.set_font(text_frame, size=font_size, color=color)

class BackgroundDesign:
    """背景设计模块"""
    @staticmethod
    def set_background_color(slide, color=(255, 255, 255)):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    @staticmethod
    def set_background_image(slide, image_path):
        background = slide.background
        fill = background.fill
        fill.type = 2  # 图片填充
        fill.user_picture(image_path)

class PPTDesigner:
    """主设计器整合所有模块"""
    def __init__(self, output_path="output.pptx", template_path=None):
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        self.output_path = output_path
        
        # 保存幻灯片尺寸，供其他模块使用
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches

    def add_cover_slide(self, title="标题", subtitle="副标题"):
        """添加封面页（使用智能文本框）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        BackgroundDesign.set_background_color(slide, (240, 240, 255))
        
        # 更新当前幻灯片尺寸
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches
        
        # 使用智能文本框添加标题
        TextBoxDesign.add_textbox(
            slide, position_profile="header", 
            text=title, level=TextLevel.TITLE
        )
        
        # 添加副标题
        TextBoxDesign.add_textbox(
            slide, left=2, top=3, width=9, height=1,
            text=subtitle, level=TextLevel.SUBTITLE
        )

    def add_content_slide(self, title="章节标题", content="正文内容", auto_adjust=True):
        """添加内容页（使用智能文本框）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 更新当前幻灯片尺寸
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches
        
        # 添加标题
        TextBoxDesign.add_textbox(
            slide, position_profile="header",
            text=title, level=TextLevel.HEADING1
        )
        
        # 添加正文内容
        TextBoxDesign.add_textbox(
            slide, position_profile="main_content",
            text=content, level=TextLevel.BODY,
            auto_adjust=auto_adjust
        )
        
        SlideHeaderFooter.set_header_footer(slide, footer_text="页脚示例")

    def add_image_slide(self, title="图片标题", image_path=None, caption="图片说明"):
        """添加图片页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 添加标题
        TextBoxDesign.add_textbox(
            slide, left=1, top=0.5, width=8, height=1,
            text=title, font_name="微软雅黑", size=Pt(36), bold=True
        )
        
        # 添加图片（如果提供了路径）
        if image_path:
            # 居中放置图片
            left = 2  # 左边距2英寸
            top = 1.5  # 上边距1.5英寸
            width = 6  # 图片宽度6英寸
            height = 4  # 图片高度4英寸
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
            
            # 添加图片说明
            if caption:
                TextBoxDesign.add_textbox(
                    slide, left=2, top=5.7, width=6, height=0.5,
                    text=caption, font_name="宋体", size=Pt(14), align=PP_ALIGN.CENTER
                )
                
        return slide

    def save(self):
        """保存PPT文件"""
        self.prs.save(self.output_path)

def test_all_features():
    """测试所有功能"""
    designer = PPTDesigner("test_output2.pptx")
    
    # 设置幻灯片尺寸（16:9）
    SlideSizeDesign.set_slide_size(designer.prs, width=Inches(13.33), height=Inches(7.5))
    
    # 添加封面页
    designer.add_cover_slide(title="清华大学PPT模板", subtitle="新雅书院 · 2023")
    
    # 添加内容页（测试超幅警告）
    designer.add_content_slide(
        title="第一章 研究方法",
        content="这里是正文内容，" * 50  # 故意制造超幅文本
    )
    
    # 添加两种艺术字
    slide = designer.prs.slides[1]
    ArtTextDesign.add_art_text(
        slide, left=1, top=4, width=3, height=1.5, 
        text="带背景艺术字", use_shape=True
    )
    ArtTextDesign.add_art_text(
        slide, left=5, top=4, width=3, height=1.5, 
        text="纯文字艺术字", use_shape=False,
        font_name="华文彩云", size=Pt(36)
    )
    
    designer.save()
    print(f"测试文件已生成: {designer.output_path}")

def test_smart_textbox():
    """测试智能文本框功能"""
    designer = PPTDesigner("smart_textbox_test.pptx")
    
    # 设置幻灯片尺寸（16:9）
    SlideSizeDesign.set_slide_size(designer.prs, width=Inches(13.33), height=Inches(7.5))
    
    # 添加封面页（智能文本框）
    designer.add_cover_slide(title="智能文本框演示", subtitle="自动调整字号与位置")
    
    # 添加内容页（测试长文本自动调整）
    long_text = "这是一段非常长的文本内容，将用于测试智能文本框的自动调整功能。" * 20
    designer.add_content_slide(
        title="自动调整文本示例",
        content=long_text
    )
    
    # 测试不同区域配置
    slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    
    # 添加标题
    TextBoxDesign.add_textbox(
        slide, position_profile="header",
        text="区域布局示例", level=TextLevel.HEADING1
    )
    
    # 左边栏
    TextBoxDesign.add_textbox(
        slide, position_profile="left_sidebar",
        text="左侧边栏内容\n• 项目一\n• 项目二\n• 项目三", 
        level=TextLevel.HEADING3
    )
    
    # 主内容区
    TextBoxDesign.add_textbox(
        slide, position_profile="main_content",
        text="这是主内容区域的文本。可以自动调整字号以适应给定空间。",
        level=TextLevel.BODY
    )
    
    # 右边栏
    TextBoxDesign.add_textbox(
        slide, position_profile="right_sidebar",
        text="右侧补充信息", 
        level=TextLevel.CAPTION
    )
    
    designer.save()
    print(f"智能文本框测试文件已生成: {designer.output_path}")

if __name__ == "__main__":
    # test_all_features()
    test_smart_textbox()