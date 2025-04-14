from pptx import Presentation
from pptx.util import Inches, Pt

# 创建一个新的演示文稿
prs = Presentation()

# 添加一个标题幻灯片（布局类型 0）
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "演示文稿标题"
subtitle.text = "副标题或你的姓名\n日期"

#一个标题和内容的幻灯片（布局类型 1）
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "第一张内容幻灯片"
content.text = "• 第一点\n• 第二点\n• 第三点"

# 添加一个带图片的幻灯片（布局类型 5）
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "带图片的幻灯片"
left = Inches(1)
top = Inches(2)
width = Inches(5)
height = Inches(3)
# 这里添加一个矩形代替图片（实际使用时可以替换为真实图片）
pic = slide.shapes.add_shape(
    1, left, top, width, height)
pic.text = "[图片位置]"

# 保存演示文稿
prs.save("presentation_template.pptx")
print("PPTX 模板已创建: presentation_template.pptx")