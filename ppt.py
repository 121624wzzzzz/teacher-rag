import os
import re
import requests
from pathlib import Path
from typing import List, Dict, Tuple
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

class LLMClient:
    """增强型大模型客户端"""
    
    def __init__(self, api_key: str, api_url: str, model_name: str):
        self.api_key = api_key
        self.api_url = api_url
        self.model_name = model_name
        self.session = requests.Session()
        self.session.headers.update({"Authorization": f"Bearer {self.api_key}"})

    def query(self, prompt: str, max_retry: int = 3) -> str:
        """带重试机制的查询函数"""
        for attempt in range(max_retry):
            try:
                response = self.session.post(
                    self.api_url,
                    json={
                        "model": self.model_name,
                        "messages": [{"role": "user", "content": prompt}]
                    },
                    timeout=300
                )
                response.raise_for_status()
                return response.json()['choices'][0]['message']['content']
            except Exception as e:
                if attempt == max_retry - 1:
                    raise RuntimeError(f"API请求失败: {str(e)}")
                continue

class PPTDesignTheme:
    """PPT主题样式配置"""
    
    THEMES = {
        "professional": {
            "title_size": Pt(36),
            "content_size": Pt(20),
            "title_color": RGBColor(12, 34, 56),
            "content_color": RGBColor(56, 78, 90),
            "background_color": RGBColor(255, 255, 255)
        },
        "tech": {
            "title_size": Pt(40),
            "content_size": Pt(24),
            "title_color": RGBColor(0, 112, 192),
            "content_color": RGBColor(47, 84, 150),
            "background_color": RGBColor(221, 235, 247)
        }
    }
    
    def __init__(self, theme_name: str = "professional"):
        self.theme = self.THEMES.get(theme_name.lower(), self.THEMES["professional"])
        
    def apply_slide_style(self, slide):
        """应用幻灯片背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["background_color"]

class SmartPPTCreator:
    """智能PPT生成器"""
    
    def __init__(self, llm_client: LLMClient, theme: str = "professional"):
        self.llm = llm_client
        self.presentation = Presentation()
        self.theme = PPTDesignTheme(theme)
        self.layouts = {
            "title": self.presentation.slide_layouts[0],
            "content": self.presentation.slide_layouts[1],
            "section": self.presentation.slide_layouts[2]
        }
        
        # 初始化默认幻灯片
        self._add_default_title_slide()

    def _add_default_title_slide(self):
        """添加默认标题页"""
        title_slide = self.presentation.slides.add_slide(self.layouts["title"])
        title_slide.shapes.title.text = "自动生成演示文稿"
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Powered by AI"
        self.theme.apply_slide_style(title_slide)

    def generate_presentation(self, topic: str, output_path: str, verbose: bool = True):
        """全流程生成演示文稿"""
        try:
            if verbose:
                print("🚀 开始生成PPT...")
            
            # 生成大纲
            outline = self._generate_valid_outline(topic, verbose)
            
            # 创建内容页
            self._generate_content_slides(outline, verbose)
            
            # 添加总结页
            self._add_summary_slide()
            
            # 保存文件
            self._safe_save(output_path, verbose)
            return output_path
            
        except Exception as e:
            self._handle_error(e)
            raise

    def _generate_valid_outline(self, topic: str, verbose: bool) -> Dict:
        """生成有效大纲（带验证和重试）"""
        MAX_RETRY = 3
        for retry in range(MAX_RETRY):
            try:
                if verbose:
                    print(f"📋 尝试生成大纲（第{retry+1}次）...")
                
                prompt = self._build_outline_prompt(topic)
                response = self.llm.query(prompt)
                outline = self._parse_markdown(response)
                
                if self._validate_outline(outline):
                    return outline
                
            except Exception as e:
                if retry == MAX_RETRY - 1:
                    raise RuntimeError("大纲生成失败，请调整提示词或检查模型响应")
                continue
                
        raise RuntimeError("无法生成有效大纲")

    def _build_outline_prompt(self, topic: str) -> str:
        """构建结构化提示词"""
        return f"""请按照以下要求生成专业PPT大纲：

主题：{topic}

格式要求：
1. 必须包含1个#开头的标题
2. 至少包含3个##开头的章节
3. 每个章节下包含2-5个-开头的要点
4. 使用纯Markdown格式，不要添加解释

示例格式：
# 主标题
## 核心概念
- 定义与特征
- 技术原理
## 应用场景
- 自然语言处理
- 计算机视觉
## 发展趋势
- 技术突破方向
- 行业应用前景
"""

    def _parse_markdown(self, text: str) -> Dict:
        """增强型Markdown解析"""
        cleaned_text = re.sub(r'```markdown|```', '', text)
        lines = [line.strip() for line in cleaned_text.split('\n') if line.strip()]
        
        outline = {"topic": "", "sections": []}
        current_section = None
        
        for line in lines:
            if line.startswith('# '):
                outline["topic"] = line[2:].strip()
            elif line.startswith('## '):
                current_section = {"title": line[3:].strip(), "content": []}
                outline["sections"].append(current_section)
            elif line.startswith('-'):
                if current_section is not None:
                    current_section["content"].append(line[1:].strip())
        
        return outline

    def _validate_outline(self, outline: Dict) -> bool:
        """大纲有效性验证"""
        return (
            len(outline.get("topic", "")) > 0 and
            len(outline.get("sections", [])) >= 2 and
            all(len(s["content"]) >= 2 for s in outline["sections"])
        )

    def _generate_content_slides(self, outline: Dict, verbose: bool):
        """生成内容幻灯片"""
        if verbose:
            print("📊 正在创建幻灯片...")
        
        # 添加主标题页
        self._add_custom_title_slide(outline["topic"])
        
        # 生成各章节内容
        for idx, section in enumerate(outline["sections"]):
            self._add_section_slide(section, idx+1, verbose)
            
            # 扩展章节内容
            expanded_content = self._expand_section_content(section)
            self._add_content_slide(section["title"], expanded_content)

    def _add_custom_title_slide(self, title: str):
        """添加自定义标题页"""
        slide = self.presentation.slides.add_slide(self.layouts["title"])
        slide.shapes.title.text = title
        self.theme.apply_slide_style(slide)

    def _add_section_slide(self, section: Dict, seq: int, verbose: bool):
        """添加章节分隔页"""
        if verbose:
            print(f"  正在生成章节：{section['title']}")
            
        slide = self.presentation.slides.add_slide(self.layouts["section"])
        title = slide.shapes.title
        title.text = f"PART {seq}\n{section['title']}"
        title.text_frame.paragraphs[0].font.color.rgb = self.theme.theme["title_color"]
        self.theme.apply_slide_style(slide)

    def _expand_section_content(self, section: Dict) -> List[str]:
        """扩展章节内容"""
        prompt = f"""请扩展以下PPT章节内容：
        
        章节标题：{section['title']}
        初始要点：{', '.join(section['content'])}
        
        要求：
        1. 生成3-5个详细子项
        2. 每个子项包含具体说明
        3. 使用层级结构（最多两级）
        4. 输出格式示例：
        - 核心概念
          - 技术定义
          - 主要特征
        - 应用场景
        """
        
        response = self.llm.query(prompt)
        return self._clean_response_content(response)

    def _clean_response_content(self, text: str) -> List[str]:
        """清理模型响应内容"""
        return [
            re.sub(r'^\s*[\d\-•*]+', '', line).strip()
            for line in text.split('\n')
            if line.strip()
        ]

    def _add_content_slide(self, title: str, content: List[str]):
        """添加内容页（带智能格式）"""
        slide = self.presentation.slides.add_slide(self.layouts["content"])
        self.theme.apply_slide_style(slide)
        
        # 设置标题
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.color.rgb = self.theme.theme["title_color"]
        
        # 添加内容
        content_box = slide.placeholders[1]
        text_frame = content_box.text_frame
        text_frame.clear()
        
        for line in content:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = self.theme.theme["content_size"]
            p.font.color.rgb = self.theme.theme["content_color"]
            p.space_after = Pt(12)

    def _add_summary_slide(self):
        """添加总结页"""
        slide = self.presentation.slides.add_slide(self.layouts["content"])
        self.theme.apply_slide_style(slide)
        
        title_box = slide.shapes.title
        title_box.text = "核心总结"
        
        content = [
            "关键技术要点回顾",
            "行业应用价值分析",
            "未来发展方向展望",
            "Q&A 交流环节"
        ]
        
        text_frame = slide.placeholders[1].text_frame
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.bold = True

    def _safe_save(self, output_path: str, verbose: bool):
        """安全保存文件"""
        try:
            self.presentation.save(output_path)
            if verbose:
                print(f"\n✅ 成功保存至：{output_path}")
        except PermissionError:
            raise RuntimeError("文件保存失败，请检查文件权限或是否已打开")
        except Exception as e:
            raise RuntimeError(f"保存失败：{str(e)}")

    def _handle_error(self, error: Exception):
        """统一错误处理"""
        print(f"\n❌ 生成失败：{str(error)}")
        print("建议检查以下内容：")
        print("1. 模型API是否可用")
        print("2. 网络连接是否正常")
        print("3. 输入主题是否明确")

if __name__ == "__main__":
    load_dotenv()
    
    # 初始化客户端
    llm_client = LLMClient(
        api_key=os.getenv("SILICONFLOW_API_KEY"),
        api_url=os.getenv("API_URL"),
        model_name=os.getenv("MODEL_NAME")
    )
    
    # 创建生成器实例
    ppt_gen = SmartPPTCreator(llm_client, theme="tech")
    
    # 用户交互
    try:
        topic = input("请输入PPT主题：").strip()
        output_path = Path.cwd() / f"{topic}_presentation.pptx"
        
        ppt_gen.generate_presentation(
            topic=topic,
            output_path=output_path,
            verbose=True
        )
        
    except KeyboardInterrupt:
        print("\n操作已取消")
    except Exception as e:
        print(f"发生未知错误：{str(e)}")