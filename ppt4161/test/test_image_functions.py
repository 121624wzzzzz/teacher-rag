# test/test_image_functions.py
import os
import unittest
from pptx import Presentation
from pptbase import PPTDesigner, BackgroundDesign

class TestImageFunctions(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        """测试前的准备工作"""
        # 确保assets目录存在
        if not os.path.exists("assets"):
            os.makedirs("assets")
        
        # 测试输出目录
        cls.output_dir = "test_outputs"
        if not os.path.exists(cls.output_dir):
            os.makedirs(cls.output_dir)

    def test_image_insertion(self):
        """测试图片插入功能"""
        # 创建测试PPT
        output_path = os.path.join(self.output_dir, "test_image_insertion.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 添加带图片的幻灯片
        slide = designer.add_image_slide(
            title="测试图片插入",
            image_path="assets/qh.jpg",  # 使用你的qh.jpg图片
            caption="清华大学测试图片"
        )
        
        # 验证图片是否被添加
        self.assertGreater(len(slide.shapes), 1)  # 至少有标题和图片
        
        # 保存PPT
        designer.save()
        print(f"图片插入测试文件已生成: {output_path}")

    def test_background_image(self):
        """测试背景图片功能"""
        # 创建测试PPT
        output_path = os.path.join(self.output_dir, "test_background_image.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 添加一个幻灯片
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        
        # 设置背景图片
        BackgroundDesign.set_background_image(slide, "assets/qh.jpg")
        
        # 添加一些文本确保可见
        TextBoxDesign.add_textbox(
            slide, left=1, top=1, width=8, height=1,
            text="带背景图片的幻灯片",
            level=TextLevel.TITLE,
            color=(255, 255, 255)  # 白色文字确保在深色背景上可见
        )
        
        # 保存PPT
        designer.save()
        print(f"背景图片测试文件已生成: {output_path}")

    def test_background_color(self):
        """测试背景颜色功能"""
        # 创建测试PPT
        output_path = os.path.join(self.output_dir, "test_background_color.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 添加一个幻灯片
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        
        # 设置背景颜色
        BackgroundDesign.set_background_color(slide, (0, 100, 200))  # 蓝色背景
        
        # 添加一些文本确保可见
        TextBoxDesign.add_textbox(
            slide, left=1, top=1, width=8, height=1,
            text="带背景颜色的幻灯片",
            level=TextLevel.TITLE,
            color=(255, 255, 255)  # 白色文字
        )
        
        # 保存PPT
        designer.save()
        print(f"背景颜色测试文件已生成: {output_path}")

if __name__ == "__main__":
    unittest.main()