# pptslides.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from .pptbase import (
    TextBoxDesign, FontDesign, BackgroundDesign,
    SlideHeaderFooter, ArtTextDesign, TextLevel
)

class SlideLayout:
    """幻灯片布局配置"""
    
    # 封面页布局
    COVER_LAYOUT = {
        "background_color": (240, 245, 255),
        "title": (0.15, 0.3, 0.7, 0.15),        # 标题居中偏上
        "subtitle": (0.2, 0.5, 0.6, 0.1),       # 副标题居中偏下
        "company": (0.6, 0.8, 0.35, 0.07)       # 公司名右下
    }
    
    # 目录页布局
    TOC_LAYOUT = {
        "background_color": (255, 255, 255),
        "title": (0.1, 0.1, 0.8, 0.12),         # 标题居顶
        "content": (0.1, 0.25, 0.4, 0.6),       # 目录左侧
        "decoration": (0.5, 0.05, 0.4, 0.25)    # 装饰元素右侧
    }
    
    # 分节页布局
    SECTION_LAYOUT = {
        "background_color": (245, 245, 250),
        "level1": (0.1, 0.1, 0.8, 0.15),        # 一级标题顶部
        "level2": (0.15, 0.35, 0.7, 0.15),      # 二级标题居中
        "level3": (0.2, 0.55, 0.6, 0.1),        # 三级标题居中偏下
        "decoration": (0.25, 0.7, 0.5, 0.02)    # 装饰线底部
    }
    
    # 内容页布局
    CONTENT_LAYOUT = {
        "background_color": (255, 255, 255),
        "title": (0.05, 0.05, 0.9, 0.12),       # 标题顶部
        "content": (0.08, 0.22, 0.85, 0.65)     # 内容区域
    }
    
    # 图文页布局
    IMAGE_CONTENT_LAYOUT = {
        "background_color": (255, 255, 255),
        "title": (0.05, 0.05, 0.9, 0.12),       # 标题顶部
        "image": (0.08, 0.22, 0.4, 0.6),        # 左侧图片
        "caption": (0.08, 0.85, 0.4, 0.08),     # 图片说明
        "content": (0.53, 0.22, 0.42, 0.65)     # 右侧文字
    }
    
    # 对比页布局
    COMPARISON_LAYOUT = {
        "background_color": (250, 250, 255),
        "title": (0.05, 0.05, 0.9, 0.12),       # 主标题顶部
        "left_title": (0.08, 0.2, 0.38, 0.1),   # 左侧标题
        "right_title": (0.54, 0.2, 0.38, 0.1),  # 右侧标题
        "left_content": (0.08, 0.32, 0.38, 0.55), # 左侧内容
        "right_content": (0.54, 0.32, 0.38, 0.55), # 右侧内容
        "separator": (0.495, 0.3, 0.01, 0.5)    # 中间分隔线
    }
    
    # 时间线页布局
    TIMELINE_LAYOUT = {
        "background_color": (252, 252, 252),
        "title": (0.05, 0.05, 0.9, 0.12),       # 标题顶部
        "timeline": (0.1, 0.45, 0.8, 0.02),     # 时间轴线
        "points": []  # 动态计算
    }
    
    # 感谢页布局
    THANK_YOU_LAYOUT = {
        "background_color": (240, 245, 255),
        "message": (0.15, 0.4, 0.7, 0.2),       # 感谢信息居中
        "contact": (0.2, 0.7, 0.6, 0.1)         # 联系方式底部
    }

class SlideDesigner:
    """PPT幻灯片结构设计模块"""
    
    def __init__(self, presentation):
        self.prs = presentation
        self.current_slide = None
        # 保存当前幻灯片尺寸，以便计算比例
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches
        
    def _inches_from_ratio(self, ratio_tuple):
        """将比例元组转换为英寸值"""
        ratio_x, ratio_y, ratio_width, ratio_height = ratio_tuple
        return (
            ratio_x * TextBoxDesign.current_slide_width,
            ratio_y * TextBoxDesign.current_slide_height,
            ratio_width * TextBoxDesign.current_slide_width,
            ratio_height * TextBoxDesign.current_slide_height
        )
        
    def add_cover_slide(self, title, subtitle=None, company=None):
        """创建封面页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.COVER_LAYOUT["background_color"]
        )
        
        # 主标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.COVER_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.TITLE,
            auto_adjust=True,
            align=PP_ALIGN.CENTER
        )
        
        # 副标题
        if subtitle:
            left, top, width, height = self._inches_from_ratio(SlideLayout.COVER_LAYOUT["subtitle"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=subtitle,
                level=TextLevel.SUBTITLE,
                align=PP_ALIGN.CENTER
            )
            
        # 公司标识（右侧底部）
        if company:
            left, top, width, height = self._inches_from_ratio(SlideLayout.COVER_LAYOUT["company"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=company,
                level=TextLevel.CAPTION,
                align=PP_ALIGN.RIGHT
            )
        
        return self.current_slide

    def add_toc_slide(self, title, items):
        """创建目录页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.TOC_LAYOUT["background_color"]
        )
        
        # 标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.TOC_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.HEADING1,
            align=PP_ALIGN.LEFT
        )
        
        # 目录内容（左侧）
        left, top, width, height = self._inches_from_ratio(SlideLayout.TOC_LAYOUT["content"])
        content_box = TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text="",
            level=TextLevel.HEADING2
        )
        
        # 添加目录条目
        for idx, item in enumerate(items):
            TextBoxDesign.add_paragraph(
                content_box.text_frame,
                text=f"{idx+1}. {item}",
                level=TextLevel.HEADING3
            )
        
        # 右侧装饰元素
        left, top, width, height = self._inches_from_ratio(SlideLayout.TOC_LAYOUT["decoration"])
        ArtTextDesign.add_art_text(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text="CONTENTS",font_size=Pt(30),shape_type="cloud",
            border_color=(100, 130, 200),  # 修改点2：建议添加边框色
            border_width=1.5,
            use_shape=True,
            shape_color=(230, 240, 255),
            text_color=(50, 80, 150)
        )
        
        return self.current_slide

    def add_section_slide(self, level1, level2=None, level3=None):
        """创建分级标题页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.SECTION_LAYOUT["background_color"]
        )
        
        # 一级标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.SECTION_LAYOUT["level1"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=level1,
            level=TextLevel.HEADING1,
            align=PP_ALIGN.LEFT
        )
        
        # 二级标题
        if level2:
            left, top, width, height = self._inches_from_ratio(SlideLayout.SECTION_LAYOUT["level2"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=level2,
                level=TextLevel.HEADING2,
                align=PP_ALIGN.CENTER
            )
            
            # 三级标题
            if level3:
                left, top, width, height = self._inches_from_ratio(SlideLayout.SECTION_LAYOUT["level3"])
                TextBoxDesign.add_textbox(
                    self.current_slide,
                    left=left, top=top, width=width, height=height,
                    text=level3,
                    level=TextLevel.HEADING3,
                    align=PP_ALIGN.CENTER
                )
        
        # 添加装饰线
        left, top, width, height = self._inches_from_ratio(SlideLayout.SECTION_LAYOUT["decoration"])
        line = self.current_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        return self.current_slide

    def add_content_slide(self, title, contents, bullet_points=True):
        """创建正文内容页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.CONTENT_LAYOUT["background_color"]
        )
        
        # 标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.CONTENT_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.HEADING3
        )
        
        # 正文内容区域
        left, top, width, height = self._inches_from_ratio(SlideLayout.CONTENT_LAYOUT["content"])
        content_box = TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text="",
            level=TextLevel.BODY,
            auto_adjust=True,
            min_size_ratio=0.7
        )
        
        # 添加内容段落
        for text in contents:
            if bullet_points:
                p = TextBoxDesign.add_paragraph(
                    content_box.text_frame,
                    text=f"• {text}",
                    level=TextLevel.BODY
                )
            else:
                p = TextBoxDesign.add_paragraph(
                    content_box.text_frame,
                    text=text,
                    level=TextLevel.BODY
                )
            p.space_after = Inches(0.1)
        
        return self.current_slide
        
    def add_image_content_slide(self, title, image_path, caption=None, text_content=None, image_transparency=0):
        """创建带图片的内容页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.IMAGE_CONTENT_LAYOUT["background_color"]
        )
        
        # 标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.IMAGE_CONTENT_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.HEADING3
        )
        
        # 图片
        left, top, width, height = self._inches_from_ratio(SlideLayout.IMAGE_CONTENT_LAYOUT["image"])
        
        # 处理图片路径
        if not os.path.exists(image_path):
            print(f"图片文件不存在: {image_path}")
            # 添加一个占位形状，表示图片缺失
            placeholder = self.current_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
            # 设置形状填充为浅灰色
            fill = placeholder.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(230, 230, 230)
            # 添加文字说明
            placeholder.text = "图片未找到"
            FontDesign.set_font(
                placeholder.text_frame, 
                level=TextLevel.CAPTION,
                align=PP_ALIGN.CENTER
            )
        else:
            # 处理透明度
            if image_transparency > 0:
                # 使用BackgroundDesign的透明度处理功能
                try:
                    BackgroundDesign.add_image_with_transparency(
                        self.current_slide, 
                        image_path,
                        left=left, 
                        top=top, 
                        width=width, 
                        height=height,
                        transparency=image_transparency
                    )
                except Exception as e:
                    print(f"设置图片透明度时出错: {e}")
                    # 出错时退回到普通图片
                    self.current_slide.shapes.add_picture(
                        image_path, Inches(left), Inches(top), 
                        Inches(width), Inches(height)
                    )
            else:
                # 不需要透明度，直接添加图片
                self.current_slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), 
                    Inches(width), Inches.height
                )
        
        # 图片说明（如果有）
        if caption:
            left, top, width, height = self._inches_from_ratio(SlideLayout.IMAGE_CONTENT_LAYOUT["caption"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=caption,
                level=TextLevel.CAPTION,
                align=PP_ALIGN.CENTER
            )
        
        # 文字内容（如果有）
        if text_content:
            left, top, width, height = self._inches_from_ratio(SlideLayout.IMAGE_CONTENT_LAYOUT["content"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=text_content,
                level=TextLevel.BODY,
                auto_adjust=True
            )
        
        return self.current_slide
    
    def add_comparison_slide(self, title, left_title, right_title, left_content, right_content):
        """创建左右对比页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.COMPARISON_LAYOUT["background_color"]
        )
        
        # 主标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.HEADING2
        )
        
        # 左侧标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["left_title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=left_title,
            level=TextLevel.HEADING3,
            align=PP_ALIGN.CENTER
        )
        
        # 右侧标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["right_title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=right_title,
            level=TextLevel.HEADING3,
            align=PP_ALIGN.CENTER
        )
        
        # 左侧内容
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["left_content"])
        left_box = TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text="",
            level=TextLevel.BODY
        )
        
        for text in left_content:
            TextBoxDesign.add_paragraph(
                left_box.text_frame,
                text=f"• {text}",
                level=TextLevel.BODY
            )
        
        # 右侧内容
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["right_content"])
        right_box = TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text="",
            level=TextLevel.BODY
        )
        
        for text in right_content:
            TextBoxDesign.add_paragraph(
                right_box.text_frame,
                text=f"• {text}",
                level=TextLevel.BODY
            )
        
        # 中间分隔线
        left, top, width, height = self._inches_from_ratio(SlideLayout.COMPARISON_LAYOUT["separator"])
        line = self.current_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        return self.current_slide
    
    def add_timeline_slide(self, title, timeline_items):
        """创建时间线页
        :param timeline_items: [(时间点, 内容描述), ...]
        """
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.TIMELINE_LAYOUT["background_color"]
        )
        
        # 标题
        left, top, width, height = self._inches_from_ratio(SlideLayout.TIMELINE_LAYOUT["title"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=title,
            level=TextLevel.HEADING2
        )
        
        # 时间线
        left, top, width, height = self._inches_from_ratio(SlideLayout.TIMELINE_LAYOUT["timeline"])
        line = self.current_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(100, 100, 200)
        
        # 添加时间点和描述
        margin = 0.1  # 左右边距比例
        usable_width = width - 2 * margin * width
        start_x = left + margin * width
        width_per_item = usable_width / len(timeline_items)
        
        for i, (time_point, description) in enumerate(timeline_items):
            point_x = start_x + i * width_per_item
            
            # 时间点标记（圆点）
            circle_size = 0.02 * TextBoxDesign.current_slide_width  # 圆点大小
            circle = self.current_slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(point_x + width_per_item/2 - circle_size/2), 
                Inches(top - circle_size/2), 
                Inches(circle_size), 
                Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(50, 100, 200)
            
            # 时间文本（圆点上方）
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=point_x, 
                top=top - 0.1 * TextBoxDesign.current_slide_height, 
                width=width_per_item, 
                height=0.06 * TextBoxDesign.current_slide_height,
                text=time_point,
                level=TextLevel.HEADING3,
                align=PP_ALIGN.CENTER
            )
            
            # 描述文本（圆点下方）
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=point_x, 
                top=top + 0.05 * TextBoxDesign.current_slide_height, 
                width=width_per_item, 
                height=0.2 * TextBoxDesign.current_slide_height,
                text=description,
                level=TextLevel.BODY,
                align=PP_ALIGN.CENTER,
                auto_adjust=True
            )
        
        return self.current_slide
    
    def add_thank_you_slide(self, text="谢谢观看！", contact_info=None):
        """创建感谢页"""
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        BackgroundDesign.set_background_color(
            self.current_slide, 
            SlideLayout.THANK_YOU_LAYOUT["background_color"]
        )
        
        # 感谢文字（大号居中）
        left, top, width, height = self._inches_from_ratio(SlideLayout.THANK_YOU_LAYOUT["message"])
        TextBoxDesign.add_textbox(
            self.current_slide,
            left=left, top=top, width=width, height=height,
            text=text,
            level=TextLevel.TITLE,
            align=PP_ALIGN.CENTER
        )
        
        # 联系信息（底部小字）
        if contact_info:
            left, top, width, height = self._inches_from_ratio(SlideLayout.THANK_YOU_LAYOUT["contact"])
            TextBoxDesign.add_textbox(
                self.current_slide,
                left=left, top=top, width=width, height=height,
                text=contact_info,
                level=TextLevel.CAPTION,
                align=PP_ALIGN.CENTER
            )
        
        return self.current_slide

# 测试样例
def test_slide_design():
    """幻灯片结构测试"""
    prs = Presentation()
    # 设置为16:9宽高比 (宽度:13.33英寸, 高度:7.5英寸)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    designer = SlideDesigner(prs)
    
    # 封面页
    designer.add_cover_slide(
        title="人工智能技术报告",
        subtitle="自然语言处理专题",
        company="智能科技公司 2023"
    )
    
    # 目录页
    designer.add_toc_slide(
        title="内容目录",
        items=[
            "技术发展现状",
            "核心算法解析",
            "应用场景分析",
            "未来发展趋势"
        ]
    )
    
    # 分级标题页
    designer.add_section_slide(
        level1="第一部分",
        level2="技术发展现状",
        level3="自然语言处理技术演进"
    )
    
    # 正文内容页
    designer.add_content_slide(
        title="关键技术解析",
        contents=[
            "Transformer架构及其优势",
            "预训练-微调范式的发展",
            "注意力机制的可解释性分析",
            "多模态融合技术的最新进展"
        ]
    )
    
    # 对比页
    designer.add_comparison_slide(
        title="传统模型 vs 深度学习模型",
        left_title="传统机器学习模型",
        right_title="深度学习模型",
        left_content=[
            "特征需要手动设计",
            "计算复杂度较低",
            "对数据量要求相对较小",
            "可解释性强"
        ],
        right_content=[
            "自动学习特征表示",
            "计算复杂度高",
            "需要大量训练数据",
            "可解释性较弱但效果优异"
        ]
    )
    
    # 时间线页
    designer.add_timeline_slide(
        title="NLP技术发展历程",
        timeline_items=[
            ("2013", "Word2Vec"),
            ("2015", "Seq2Seq"),
            ("2017", "Transformer"),
            ("2018", "BERT"),
            ("2020", "GPT-3")
        ]
    )
    
    # 带图片的内容页 - 假设有图片路径
    try:
        designer.add_image_content_slide(
            title="多模态模型架构",
            image_path="example.jpg",  # 假设存在此图片
            caption="CLIP模型架构示意图",
            text_content="多模态模型能够同时处理文本和图像信息，实现跨模态理解和生成。CLIP模型通过对比学习，将文本和图像映射到同一语义空间，从而实现零样本迁移能力。"
        )
    except FileNotFoundError:
        print("示例图片不存在，跳过图片页创建")
    
    # 带图片的内容页 - 使用透明度
    try:
        designer.add_image_content_slide(
            title="带透明度效果的多模态模型架构",
            image_path="assets/qh.jpg",  # 使用已有的图片
            caption="透明度设置为0.3的演示图",
            text_content="这是一个演示图片透明度效果的例子。通过调整透明度参数，可以使图片呈现出不同的视觉效果，更好地与内容融合。",
            image_transparency=0.3  # 设置30%透明度
        )
    except FileNotFoundError:
        print("图片不存在，跳过图片页创建")
    
    # 感谢页
    designer.add_thank_you_slide(
        text="感谢聆听！",
        contact_info="联系方式：example@ai.com | 网站：www.example.ai"
    )
    
    # 保存测试文件
    prs.save("slide_design_test.pptx")
    print("测试文件已生成：slide_design_test.pptx")

def test_presentation_workflow():
    """测试完整演示文稿工作流"""
    prs = Presentation()
    # 设置为16:9宽高比 (宽度:13.33英寸, 高度:7.5英寸)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    designer = SlideDesigner(prs)
    
    # 1. 封面页 - 添加背景图片
    cover_slide = designer.add_cover_slide(
        title="人工智能发展趋势分析",
        subtitle="从大语言模型看AI的未来",
        company="AI研究院 2023"
    )
    
    # 设置封面背景图片
    BackgroundDesign.set_background_image(cover_slide, "assets/qh.jpg", transparency=0.9)
    
    # 2. 目录页
    designer.add_toc_slide(
        title="报告大纲",
        items=[
            "AI历史发展回顾",
            "大语言模型的技术突破",
            "当前应用场景分析",
            "面临的挑战与局限",
            "未来发展预测"
        ]
    )
    
    # 3. 第一部分：历史发展
    section_slide = designer.add_section_slide(
        level1="第一部分",
        level2="AI历史发展回顾",
        level3="从规则系统到深度学习"
    )
    
    # 设置分节页背景图片
    BackgroundDesign.set_background_image(section_slide, "assets/qh.jpg", transparency=0.6)
    
    designer.add_timeline_slide(
        title="AI关键发展节点",
        timeline_items=[
            ("1950s", "图灵测试提出"),
            ("1980s", "专家系统流行"),
            ("2000s", "统计机器学习"),
            ("2012", "深度学习兴起"),
            ("2022", "大语言模型爆发")
        ]
    )
    
    # 4. 第二部分：技术突破
    designer.add_section_slide(
        level1="第二部分",
        level2="大语言模型的技术突破"
    )
    
    designer.add_content_slide(
        title="核心技术创新",
        contents=[
            "预训练自然语言模型的规模化",
            "Transformer架构的改进与优化",
            "自监督学习的广泛应用",
            "强化学习和人类反馈的引入",
            "多模态能力的整合"
        ]
    )
    
    # 5. 第三部分：应用场景
    designer.add_section_slide(
        level1="第三部分",
        level2="当前应用场景分析"
    )
    
    designer.add_content_slide(
        title="主要应用领域",
        contents=[
            "内容创作与辅助",
            "编程与软件开发",
            "教育辅助和个性化学习",
            "客户服务与智能问答",
            "专业领域辅助决策"
        ]
    )
    
    # 6. 第四部分：挑战与局限
    designer.add_section_slide(
        level1="第四部分",
        level2="面临的挑战与局限"
    )
    
    designer.add_comparison_slide(
        title="大语言模型的优势与局限",
        left_title="主要优势",
        right_title="关键局限",
        left_content=[
            "强大的语言理解能力",
            "零样本学习和少样本学习",
            "多任务适应性",
            "持续进化的潜力"
        ],
        right_content=[
            "事实准确性难以保证",
            "存在偏见和伦理问题",
            "推理能力有限",
            "环境适应性不足",
            "计算资源需求大"
        ]
    )
    
    # 7. 第五部分：未来预测
    designer.add_section_slide(
        level1="第五部分",
        level2="未来发展预测"
    )
    
    designer.add_content_slide(
        title="未来五年可能的发展方向",
        contents=[
            "多模态理解与生成的进一步融合",
            "具备更强规划和推理能力的模型",
            "领域专精模型与通用模型的协同",
            "适应性和持续学习能力的增强",
            "更高效的训练和推理范式"
        ]
    )
    
    # 8. 感谢页 - 添加背景图片
    thank_slide = designer.add_thank_you_slide(
        text="感谢您的关注！",
        contact_info="联系我们：research@aiinstitute.org"
    )
    
    # 设置感谢页背景图片
    BackgroundDesign.set_background_image(thank_slide, "assets/qh.jpg", transparency=0.9)
    
    # 保存演示文稿
    prs.save("ai_trends_presentation.pptx")
    print("完整演示文稿已生成：ai_trends_presentation.pptx")

if __name__ == "__main__":
    test_slide_design()
    # 取消下面的注释来运行完整演示文稿示例
    test_presentation_workflow()