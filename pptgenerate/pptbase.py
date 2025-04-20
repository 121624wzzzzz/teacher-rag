#  pptbase.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random, os
from enum import Enum, auto
from typing import Optional, Literal
# 导入Pillow相关库
from PIL import Image
import io

class TextLevel(Enum):
    """文本层级枚举"""
    TITLE = auto()      # 主标题
    SUBTITLE = auto()   # 副标题
    HEADING1 = auto()   # 一级标题
    HEADING2 = auto()   # 二级标题
    HEADING3 = auto()   # 三级标题
    BODY = auto()       # 正文
    CAPTION = auto()    # 说明文字

# 默认字号配置
DEFAULT_FONT_SIZES = {
    TextLevel.TITLE: Pt(36),
    TextLevel.SUBTITLE: Pt(24),
    TextLevel.HEADING1: Pt(28),
    TextLevel.HEADING2: Pt(22),
    TextLevel.HEADING3: Pt(18),
    TextLevel.BODY: Pt(14),
    TextLevel.CAPTION: Pt(12)
}

class FontDesign:
    """增强版字体设计模块"""
    
    # 字体层级配置
    LEVEL_STYLES = {
        TextLevel.TITLE: {
            "font_name": "微软雅黑",
            "color": (0, 0, 0),
            "bold": True,
            "align": PP_ALIGN.CENTER
        },
        TextLevel.SUBTITLE: {
            "font_name": "微软雅黑",
            "color": (100, 100, 100),
            "bold": False,
            "align": PP_ALIGN.CENTER
        },
        TextLevel.HEADING1: {
            "font_name": "微软雅黑",
            "color": (0, 0, 139),  # 深蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.HEADING2: {
            "font_name": "微软雅黑",
            "color": (0, 70, 130),  # 中蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.HEADING3: {
            "font_name": "微软雅黑",
            "color": (0, 100, 120),  # 淡蓝色
            "bold": True,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.BODY: {
            "font_name": "宋体",
            "color": (0, 0, 0),
            "bold": False,
            "align": PP_ALIGN.LEFT
        },
        TextLevel.CAPTION: {
            "font_name": "宋体",
            "color": (100, 100, 100),
            "bold": False,
            "align": PP_ALIGN.CENTER
        }
    }
    
    @staticmethod
    def set_font(text_frame, level=None, font_name=None, size=None, color=None, bold=None, align=None):
        """
        设置字体样式
        :param level: 文本层级
        :param font_name, size, color, bold, align: 可选参数，覆盖默认设置
        """
        paragraph = text_frame.paragraphs[0]
        font = paragraph.font
        
        # 如果指定了level，使用预定义样式
        if level is not None:
            style = FontDesign.LEVEL_STYLES.get(level, {})
            font.name = font_name or style.get("font_name", "微软雅黑")
            font.size = size or DEFAULT_FONT_SIZES.get(level, Pt(12))
            font.bold = bold if bold is not None else style.get("bold", False)
            font.color.rgb = RGBColor(*(color or style.get("color", (0, 0, 0))))
            paragraph.alignment = align or style.get("align", PP_ALIGN.LEFT)
        else:
            # 兼容旧接口
            font.name = font_name or "微软雅黑"
            font.size = size or Pt(12)
            font.bold = bold if bold is not None else False
            font.color.rgb = RGBColor(*(color or (0, 0, 0)))
            paragraph.alignment = align or PP_ALIGN.LEFT

class TextBoxDesign:
    """智能文本框设计模块"""
    
    # 静态类属性，存储当前幻灯片的尺寸
    current_slide_width = 13.33  # 默认16:9宽度（英寸）
    current_slide_height = 7.5   # 默认16:9高度（英寸）
    
    # 预定义的位置配置（以幻灯片尺寸的比例表示）
    POSITION_PROFILES = {
        "header": (0.2, 0.1, 0.6, 0.15),        # 页眉区域
        "left_sidebar": (0.05, 0.3, 0.25, 0.6), # 左侧边栏
        "right_sidebar": (0.7, 0.3, 0.25, 0.6), # 右侧边栏
        "main_content": (0.3, 0.3, 0.35, 0.6),  # 主内容区
        "footer": (0.2, 0.85, 0.6, 0.1)         # 页脚区域
    }
    
    @staticmethod
    def _estimate_lines(text, font_size, box_width, font_name='微软雅黑'):
        """精确估算文本行数（考虑字体和段落）"""
        # 中文字体列表
        chinese_fonts = ['微软雅黑', '宋体', 'SimSun', 'Microsoft YaHei', '黑体', 'SimHei', '楷体', 'KaiTi']
        is_chinese = any(f in font_name for f in chinese_fonts)
        
        # 根据字体类型确定字符宽度系数
        # 中文字符通常是方形的，宽度接近高度
        char_width_factor = 1.0 if is_chinese else 0.6
        char_width = font_size.pt * char_width_factor
        
        if char_width <= 0:
            return len(text.split('\n'))  # 防止除零错误
        
        # 计算每行字符数
        box_width_pt = box_width * 72  # 将英寸转换为点
        chars_per_line = max(1, int(box_width_pt / char_width))
        
        # 处理段落换行
        total_lines = 0
        for paragraph in text.split('\n'):
            if not paragraph:  # 空段落
                total_lines += 1
                continue
                
            # 字符数 / 每行字符 = 行数（向上取整）
            para_lines = (len(paragraph) + chars_per_line - 1) // chars_per_line
            total_lines += max(1, para_lines)  # 至少占一行
            
        return total_lines

    @staticmethod
    def add_textbox(slide, left=None, top=None, width=None, height=None, 
                    text="", level=None, position_profile=None, 
                    auto_adjust=True, min_size_ratio=0.6, word_wrap=True, **font_kwargs):
        """
        添加智能文本框
        :param slide: 幻灯片对象
        :param left, top, width, height: 位置和尺寸（英寸）
        :param text: 文本内容
        :param level: 文本层级
        :param position_profile: 位置配置名称
        :param auto_adjust: 是否自动调整字号
        :param min_size_ratio: 最小可接受字号比例(相对于默认字号)
        :param word_wrap: 是否启用自动换行
        :param font_kwargs: 其他字体参数
        """
        # 处理位置配置
        if position_profile:
            try:
                # 使用类属性存储的幻灯片尺寸
                slide_width = TextBoxDesign.current_slide_width
                slide_height = TextBoxDesign.current_slide_height
                
                # 计算实际位置和尺寸
                profile = TextBoxDesign.POSITION_PROFILES.get(position_profile)
                if profile:
                    left_ratio, top_ratio, width_ratio, height_ratio = profile
                    left = left or (slide_width * left_ratio)
                    top = top or (slide_height * top_ratio)
                    width = width or (slide_width * width_ratio)
                    height = height or (slide_height * height_ratio)
            except Exception as e:
                print(f"警告：无法应用位置配置'{position_profile}'，使用默认位置。错误：{str(e)}")
        
        # 确保所有位置参数存在
        if None in (left, top, width, height):
            # 默认值
            left = left or 1
            top = top or 1
            width = width or 8
            height = height or 1
        
        # 创建文本框 - 使用全局导入的Inches
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        
        # 设置自动换行属性
        text_frame.word_wrap = word_wrap
        
        text_frame.text = text
        
        # 获取默认字号和字体信息
        if level:
            style = FontDesign.LEVEL_STYLES.get(level, {})
            default_size = DEFAULT_FONT_SIZES.get(level, Pt(12))
            current_font = font_kwargs.get('font_name', style.get('font_name', '微软雅黑'))
        else:
            default_size = font_kwargs.get('size', Pt(12))
            current_font = font_kwargs.get('font_name', '微软雅黑')
        
        # 设置默认字体
        if level:
            FontDesign.set_font(text_frame, level=level, **font_kwargs)
        else:
            FontDesign.set_font(text_frame, **font_kwargs)
        
        # 自动调整字号（优化版本，使用二分查找）
        if auto_adjust and text:  # 修正 && 为 and
            min_size = Pt(max(8, default_size.pt * min_size_ratio))  # 最小不低于8pt
            
            # 二分查找参数
            low, high = min_size.pt, default_size.pt
            best_size = min_size.pt  # 初始设为最小值
            
            while low <= high:
                mid = (low + high) // 2
                test_size = Pt(mid)
                
                # 计算预估行数和最大行数
                estimated = TextBoxDesign._estimate_lines(text, test_size, width, current_font)
                line_height = test_size.pt * 1.2  # 1.2倍行高
                max_lines = int((height * 72) / line_height)  # 总高度能容纳的行数
                
                if estimated <= max_lines:
                    # 当前字号可行，记录并尝试更大字号
                    best_size = mid
                    low = mid + 1
                else:
                    # 当前字号太大，尝试更小字号
                    high = mid - 1
            
            # 应用最佳字号
            final_size = Pt(best_size)
            
            # 如果字号有调整，更新字体
            if (final_size.pt != default_size.pt):
                print(f"信息：文本框（位置：{left},{top}）字号从{default_size.pt}pt自动调整为{final_size.pt}pt")
                if level:
                    FontDesign.set_font(text_frame, level=level, size=final_size, **font_kwargs)
                else:
                    FontDesign.set_font(text_frame, size=final_size, **font_kwargs)
            
            # 最终验证是否仍然超幅
            final_estimate = TextBoxDesign._estimate_lines(text, final_size, width, current_font)
            final_max = int((height * 72) / (final_size.pt * 1.2))
            if final_estimate > final_max:
                print(f"警告：文本框可能超幅（位置：{left},{top}，预估行数：{final_estimate} > 最大行数：{final_max}）")
        
        return textbox

    @staticmethod
    def add_paragraph(text_frame, text, level=None, **font_kwargs):
        """
        添加新段落并保持样式一致
        :param text_frame: 文本框架
        :param text: 段落文本
        :param level: 文本层级
        :param font_kwargs: 字体参数
        """
        # 确保文本框启用自动换行
        text_frame.word_wrap = True
        
        paragraph = text_frame.add_paragraph()
        paragraph.text = text
        
        if level:
            font_style = FontDesign.LEVEL_STYLES.get(level, {})
            for idx, run in enumerate(paragraph.runs):
                run.font.name = font_kwargs.get('font_name', font_style.get('font_name', '微软雅黑'))
                run.font.size = font_kwargs.get('size', DEFAULT_FONT_SIZES.get(level, Pt(12)))
                run.font.bold = font_kwargs.get('bold', font_style.get('bold', False))
                run.font.color.rgb = RGBColor(*font_kwargs.get('color', font_style.get('color', (0, 0, 0))))
            paragraph.alignment = font_kwargs.get('align', font_style.get('align', PP_ALIGN.LEFT))
        else:
            # 使用传入的参数
            for idx, run in enumerate(paragraph.runs):
                run.font.name = font_kwargs.get('font_name', '微软雅黑')
                run.font.size = font_kwargs.get('size', Pt(12))
                run.font.bold = font_kwargs.get('bold', False)
                run.font.color.rgb = RGBColor(*font_kwargs.get('color', (0, 0, 0)))
            paragraph.alignment = font_kwargs.get('align', PP_ALIGN.LEFT)
        
        return paragraph

class SlideHeaderFooter:
    """页眉页脚设计模块"""
    @staticmethod
    def set_header_footer(slide, header_text="", footer_text="", footer_font_size=Pt(10)):
        # 页眉（需通过修改XML实现）
        if header_text:
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 12:  # 12对应页眉
                    shape.text = header_text

        # 页脚（简单实现）
        if footer_text:
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8.0), Inches(0.5))
            footer.text = footer_text
            FontDesign.set_font(footer.text_frame, size=footer_font_size, align=PP_ALIGN.CENTER)

class SlideSizeDesign:
    """幻灯片尺寸设计模块"""
    @staticmethod
    def set_slide_size(presentation, width=Inches(13.33), height=Inches(7.5)):
        presentation.slide_width = width
        presentation.slide_height = height

class ArtTextDesign:
    """支持多框型选择的艺术字设计模块"""
    
    # 预定义的框型映射表
    SHAPE_MAP = {
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "cloud": MSO_SHAPE.CLOUD,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "circle": MSO_SHAPE.OVAL,
        "ribbon": MSO_SHAPE.CURVED_UP_RIBBON,
        "star": MSO_SHAPE.STAR_10_POINT,
        # 可继续扩展其他形状...
    }

    @staticmethod
    def add_art_text(
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str = "ART",
        font_size: int = Pt(36),
        # 新版参数（兼容旧版use_shape）
        shape_type: Optional[Literal["rounded_rect", "cloud", "arrow", "circle", "ribbon", "star"]] = None,
        use_shape: bool = False,  # 保留旧版参数
        shape_color: tuple = (240, 240, 240),
        border_color: tuple = (150, 150, 150),
        border_width: float = 1.0,
        text_color: Optional[tuple] = None,
        **font_kwargs
    ):
        """
        添加艺术字（新版支持多框型）
        
        参数：
        :param shape_type: 框型类型，可选："rounded_rect"|"cloud"|"arrow"|"circle"|"ribbon"|"star"|None
        :param use_shape: 旧版参数（True等效于shape_type="rounded_rect"）
        :param border_color: 边框颜色RGB值
        :param border_width: 边框宽度（磅）
        其他参数与原版一致
        """
        # 兼容旧版use_shape参数
        if use_shape and shape_type is None:
            shape_type = "rounded_rect"

        # 创建形状或文本框
        if shape_type is not None:
            # 验证形状类型有效性
            if shape_type not in ArtTextDesign.SHAPE_MAP:
                raise ValueError(f"不支持的shape_type: {shape_type}，可选: {list(ArtTextDesign.SHAPE_MAP.keys())}")
            
            # 带指定框型的艺术字
            shape = slide.shapes.add_shape(
                ArtTextDesign.SHAPE_MAP[shape_type],
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            shape.text = text
            
            # 设置形状样式
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*shape_color)
            
            line = shape.line
            line.color.rgb = RGBColor(*border_color)
            line.width = Pt(border_width)
        else:
            # 纯文字模式
            shape = TextBoxDesign.add_textbox(
                slide, left, top, width, height, text,
                check_overflow=False,
                **font_kwargs
            )
        
        # 统一设置文字样式
        text_frame = shape.text_frame
        # 确保艺术字文本框也启用自动换行
        text_frame.word_wrap = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        color = text_color or (
            random.randint(0, 200),  # 限制亮度避免浅色背景看不清
            random.randint(0, 200),
            random.randint(0, 200)
        )
        FontDesign.set_font(text_frame, size=font_size, color=color)

class BackgroundDesign:
    """背景设计模块"""
    @staticmethod
    def set_background_color(slide, color=(255, 255, 255)):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    @staticmethod
    def adjust_image_transparency(image_path, transparency):
        """
        使用Pillow调整图片透明度
        
        参数:
            image_path: 图片路径
            transparency: 透明度(0-1)，0为完全不透明，1为完全透明
        
        返回:
            BytesIO对象，包含处理后的PNG图片数据
        """
        img = Image.open(image_path)
        
        # 转换为RGBA模式（如果还不是）
        if img.mode != 'RGBA':
            img = img.convert('RGBA')
        
        # 获取图片数据
        data = img.getdata()
        
        # 创建新数据，调整透明度
        new_data = []
        alpha = int(255 * (1 - transparency))  # 计算alpha值
        
        for item in data:
            # 保持RGB不变，只调整alpha通道
            new_data.append((item[0], item[1], item[2], alpha))
        
        # 应用新数据
        img.putdata(new_data)
        
        # 将图片保存到内存中的BytesIO对象
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)  # 重置指针到开头
        
        return img_bytes

    @staticmethod
    def add_image_with_transparency(slide, image_path, left=0, top=0, width=None, height=None, transparency=0.0, as_background=False):
        """
        添加带透明度的图片（可用于背景或普通图片）
        
        参数:
            slide: 幻灯片对象
            image_path: 图片路径
            left, top: 位置（英寸）
            width, height: 尺寸（英寸），None表示使用原始尺寸
            transparency: 透明度（0-1）
            as_background: 是否作为背景（放置在最底层）
        
        返回:
            添加的图片对象
        """
        # 获取幻灯片尺寸，用于默认值
        slide_width = TextBoxDesign.current_slide_width
        slide_height = TextBoxDesign.current_slide_height
        
        # 如果作为背景，默认使用整个幻灯片尺寸
        if as_background:
            width = width or slide_width
            height = height or slide_height
        
        # 处理透明度
        if transparency > 0:
            # 使用Pillow处理透明度
            img_bytes = BackgroundDesign.adjust_image_transparency(image_path, transparency)
            
            # 添加处理后的图片
            pic = slide.shapes.add_picture(
                img_bytes, 
                Inches(left), Inches(top),
                Inches(width) if width else None,
                Inches(height) if height else None
            )
        else:
            # 不需要透明度，直接添加
            pic = slide.shapes.add_picture(
                image_path, 
                Inches(left), Inches(top),
                Inches(width) if width else None,
                Inches(height) if height else None
            )
            
        # 如果作为背景，移动到最底层
        if as_background:
            pic_element = pic._element
            spTree = slide.shapes._spTree
            spTree.remove(pic_element)
            spTree.insert(0, pic_element)
            
        return pic

    @staticmethod
    def set_background_image(slide, image_path, transparency=0.0):
        """设置幻灯片背景图片
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            transparency: 透明度（0-1）
        """
        # 添加背景图片
        pic = BackgroundDesign.add_image_with_transparency(
            slide, image_path, 
            left=0, top=0, 
            width=TextBoxDesign.current_slide_width,
            height=TextBoxDesign.current_slide_height,
            transparency=transparency,
            as_background=True
        )
        
        return pic

    @staticmethod
    def set_shape_background_image(shape, image_path):
        """设置形状的背景为图片
        
        Args:
            shape: 要设置背景的形状对象
            image_path: 图片文件路径
        """
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"图片文件不存在: {image_path}")
            
        try:
            # 设置形状填充为图片
            fill = shape.fill
            
            # 首先确保填充类型是正确的
            # 注意：python-pptx API中实际调用顺序很重要
            fill.solid()  # 先设置为纯色填充
            
            # 尝试直接设置用户图片（可能在某些版本中工作）
            try:
                fill.user_picture(image_path)
            except Exception:
                # 如果直接设置失败，尝试以下替代方法
                try:
                    # 方法1: 尝试使用图案填充方法
                    fill.patterned()
                    fill.pattern = 1  # 设置为实心图案
                    fill.user_picture(image_path)
                except Exception:
                    # 方法2: 最后尝试使用渐变填充方法
                    fill.gradient()
                    fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)
                    fill.user_picture(image_path)
        except Exception as e:
            print(f"设置形状背景图片时出错: {e}")
            
            # 退回到纯色填充
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(200, 200, 200)  # 浅灰色
            
            # 添加错误提示文本
            if hasattr(shape, 'text_frame'):
                shape.text = "图片加载失败"
                if hasattr(shape.text_frame, 'paragraphs') and len(shape.text_frame.paragraphs) > 0:
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

class PPTDesigner:
    """主设计器整合所有模块"""
    def __init__(self, output_path="output.pptx", template_path=None):
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        self.output_path = output_path
        
        # 保存幻灯片尺寸，供其他模块使用
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches

    def add_cover_slide(self, title="标题", subtitle="副标题"):
        """添加封面页（使用智能文本框）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        BackgroundDesign.set_background_color(slide, (240, 240, 255))
        
        # 更新当前幻灯片尺寸
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches
        
        # 使用智能文本框添加标题
        TextBoxDesign.add_textbox(
            slide, position_profile="header", 
            text=title, level=TextLevel.TITLE
        )
        
        # 添加副标题
        TextBoxDesign.add_textbox(
            slide, left=2, top=3, width=9, height=1,
            text=subtitle, level=TextLevel.SUBTITLE
        )

    def add_content_slide(self, title="章节标题", content="正文内容", auto_adjust=True):
        """添加内容页（使用智能文本框）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 更新当前幻灯片尺寸
        TextBoxDesign.current_slide_width = self.prs.slide_width.inches
        TextBoxDesign.current_slide_height = self.prs.slide_height.inches
        
        # 添加标题
        TextBoxDesign.add_textbox(
            slide, position_profile="header",
            text=title, level=TextLevel.HEADING1
        )
        
        # 添加正文内容，确保auto_adjust为True并设置较小的min_size_ratio
        TextBoxDesign.add_textbox(
            slide, position_profile="main_content",
            text=content, level=TextLevel.BODY,
            auto_adjust=True, min_size_ratio=0.4
        )
        
        SlideHeaderFooter.set_header_footer(slide, footer_text="页脚示例")

    def add_image_slide(self, title="图片标题", image_path=None, caption="图片说明"):
        """添加图片页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 添加标题
        TextBoxDesign.add_textbox(
            slide, left=1, top=0.5, width=8, height=1,
            text=title, font_name="微软雅黑", size=Pt(36), bold=True
        )
        
        # 添加图片（如果提供了路径）
        if image_path:
            # 居中放置图片
            left = 2  # 左边距2英寸
            top = 1.5  # 上边距1.5英寸
            width = 6  # 图片宽度6英寸
            height = 4  # 图片高度4英寸
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
            
            # 添加图片说明
            if caption:
                TextBoxDesign.add_textbox(
                    slide, left=2, top=5.7, width=6, height=0.5,
                    text=caption, font_name="宋体", size=Pt(14), align=PP_ALIGN.CENTER
                )
                
        return slide

    def save(self):
        """保存PPT文件"""
        self.prs.save(self.output_path)

def test_all_features():
    """测试所有功能"""
    designer = PPTDesigner("test_output2.pptx")
    
    # 设置幻灯片尺寸（16:9）
    SlideSizeDesign.set_slide_size(designer.prs, width=Inches(13.33), height=Inches(7.5))
    
    # 添加封面页
    designer.add_cover_slide(title="清华大学PPT模板", subtitle="新雅书院 · 2023")
    
    # 添加内容页（测试超幅警告）
    designer.add_content_slide(
        title="第一章 研究方法",
        content="这里是正文内容，" * 50  # 故意制造超幅文本
    )
    
    # 添加两种艺术字
    slide = designer.prs.slides[1]
    ArtTextDesign.add_art_text(
        slide, left=1, top=4, width=3, height=1.5, 
        text="带背景艺术字", use_shape=True
    )
    ArtTextDesign.add_art_text(
        slide, left=5, top=4, width=3, height=1.5, 
        text="纯文字艺术字", use_shape=False,
        font_name="华文彩云", size=Pt(36)
    )
    
    designer.save()
    print(f"测试文件已生成: {designer.output_path}")

def test_smart_textbox():
    """测试智能文本框功能"""
    designer = PPTDesigner("smart_textbox_test.pptx")
    
    # 设置幻灯片尺寸（16:9）
    SlideSizeDesign.set_slide_size(designer.prs, width=Inches(13.33), height=Inches(7.5))
    
    # 添加封面页（智能文本框）
    designer.add_cover_slide(title="智能文本框演示", subtitle="自动调整字号与位置")
    
    # 添加内容页（测试长文本自动调整）- 使用更长的文本确保触发调整
    long_text = "这是一段非常长的文本内容，将用于测试智能文本框的自动调整功能。" * 30
    designer.add_content_slide(
        title="自动调整文本示例",
        content=long_text
    )
    
    # 创建另一个测试页面，直接使用TextBoxDesign.add_textbox来测试
    slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    
    # 添加标题
    TextBoxDesign.add_textbox(
        slide, position_profile="header",
        text="自动调整字号示例", level=TextLevel.HEADING1
    )
    
    # 直接添加一个超长文本，确保触发字号调整
    TextBoxDesign.add_textbox(
        slide, left=1, top=2, width=6, height=3,
        text="这是另一段测试文本，用于验证字号自动调整功能是否正常工作。" * 25, 
        level=TextLevel.BODY,
        auto_adjust=True,
        min_size_ratio=0.4  # 允许字号缩小到默认的40%
    )
    
    # 测试区域配置的幻灯片
    slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    
    # 添加标题
    TextBoxDesign.add_textbox(
        slide, position_profile="header",
        text="区域布局示例", level=TextLevel.HEADING1
    )
    
    # 左边栏
    TextBoxDesign.add_textbox(
        slide, position_profile="left_sidebar",
        text="左侧边栏内容\n• 项目一\n• 项目二\n• 项目三", 
        level=TextLevel.HEADING3
    )
    
    # 主内容区
    TextBoxDesign.add_textbox(
        slide, position_profile="main_content",
        text="这是主内容区域的文本。可以自动调整字号以适应给定空间。这是主内容区域的文本。可以自动调整字号以适应给定空间。这是主内容区域的文本。可以自动调整字号以适应给定空间。",
        level=TextLevel.BODY
    )
    
    # 右边栏
    TextBoxDesign.add_textbox(
        slide, position_profile="right_sidebar",
        text="右侧补充信息", 
        level=TextLevel.CAPTION
    )
    
    designer.save()
    print(f"智能文本框测试文件已生成: {designer.output_path}")

def test_background_transparency():
    """测试背景透明度功能"""
    designer = PPTDesigner("background_transparency_test.pptx")
    
    # 设置幻灯片尺寸（16:9）
    SlideSizeDesign.set_slide_size(designer.prs, width=Inches(13.33), height=Inches(7.5))
    
    # 添加封面页（使用智能文本框）
    designer.add_cover_slide(title="图片透明度演示", subtitle="基于Pillow的透明度控制")
    
    # 添加背景图片透明度测试
    slide1 = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    BackgroundDesign.set_background_image(slide1, "assets/qh.jpg", transparency=0.3)
    TextBoxDesign.add_textbox(
        slide1, left=1, top=1, width=5, height=1,
        text="背景图片 - 透明度0.3", level=TextLevel.HEADING2
    )
    
    # 添加不同透明度的背景图片
    slide2 = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    BackgroundDesign.set_background_image(slide2, "assets/qh.jpg", transparency=0.6)
    TextBoxDesign.add_textbox(
        slide2, left=1, top=1, width=5, height=1,
        text="背景图片 - 透明度0.6", level=TextLevel.HEADING2
    )
    
    # 添加普通图片透明度测试
    slide3 = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
    TextBoxDesign.add_textbox(
        slide3, position_profile="header",
        text="普通图片透明度测试", level=TextLevel.HEADING1
    )
    
    # 添加不同透明度的普通图片
    transparencies = [0.0, 0.3, 0.6, 0.8]
    for i, trans in enumerate(transparencies):
        BackgroundDesign.add_image_with_transparency(
            slide3, "assets/qh.jpg",
            left=2 + (i*2.5), top=2,
            width=2, height=1.5,
            transparency=trans
        )
        
        TextBoxDesign.add_textbox(
            slide3, left=2 + (i*2.5), top=3.7, width=2, height=0.5,
            text=f"透明度: {trans}", level=TextLevel.CAPTION
        )
    
    designer.save()
    print(f"图片透明度测试文件已生成: {designer.output_path}")

if __name__ == "__main__":
    # test_all_features()
    test_smart_textbox()
    test_background_transparency()