import os
import sys
import unittest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 添加项目根目录到Python路径
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

# 现在可以导入pptbase模块
from pptbase import PPTDesigner, BackgroundDesign, ArtTextDesign, TextBoxDesign, TextLevel, SlideSizeDesign

class TestImageBase(unittest.TestCase):
    """测试图片基础功能"""
    
    @classmethod
    def setUpClass(cls):
        """测试前的准备工作"""
        # 确保assets目录存在
        if not os.path.exists("assets"):
            os.makedirs("assets")
        
        # 测试输出目录
        cls.output_dir = "test_outputs/image_base"
        if not os.path.exists(cls.output_dir):
            os.makedirs(cls.output_dir)
    
    def test_basic_image_insertion(self):
        """测试基本图片插入功能"""
        output_path = os.path.join(self.output_dir, "basic_image_insertion.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 设置幻灯片尺寸（16:9）
        SlideSizeDesign.set_slide_size(designer.prs)
        
        # 添加标题幻灯片
        designer.add_cover_slide(title="图片功能测试", subtitle="测试PPT基础图片功能")
        
        # 添加带图片的幻灯片
        slide = designer.add_image_slide(
            title="清华大学图片测试",
            image_path="assets/qh.jpg",
            caption="清华大学校园风光"
        )
        
        # 验证图片是否被添加
        self.assertGreater(len(slide.shapes), 2)  # 至少有标题、图片和说明文字
        
        designer.save()
        print(f"基本图片插入测试文件已生成: {output_path}")
    
    def test_multiple_images(self):
        """测试多图片布局"""
        output_path = os.path.join(self.output_dir, "multiple_images.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 添加一个空白幻灯片
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        
        # 添加标题
        TextBoxDesign.add_textbox(
            slide, left=1, top=0.5, width=8, height=1,
            text="多图片布局测试", level=TextLevel.TITLE
        )
        
        # 左侧图片
        left_pic = slide.shapes.add_picture("assets/qh.jpg", 
                                           Inches(1), Inches(2), 
                                           Inches(5), Inches(3))
        
        # 右侧小图片（同一张图片用不同尺寸）
        right_pic = slide.shapes.add_picture("assets/qh.jpg", 
                                            Inches(7), Inches(2), 
                                            Inches(3), Inches(2))
        
        # 添加说明文字
        TextBoxDesign.add_textbox(
            slide, left=7, top=4.2, width=3, height=0.5,
            text="缩略图", level=TextLevel.CAPTION
        )
        
        designer.save()
        print(f"多图片布局测试文件已生成: {output_path}")
    
    def test_background_variations(self):
        """测试背景图片和颜色变体"""
        output_path = os.path.join(self.output_dir, "background_variations.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 创建带颜色背景的幻灯片
        for color, name in [
            ((100, 100, 255), "淡蓝色背景"),
            ((100, 255, 100), "淡绿色背景"),
            ((255, 100, 100), "淡红色背景")
        ]:
            slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
            
            # 设置背景颜色
            BackgroundDesign.set_background_color(slide, color)
            
            # 添加标题
            TextBoxDesign.add_textbox(
                slide, left=1, top=1, width=10, height=1,
                text=name, level=TextLevel.HEADING1
            )
            
            # 添加艺术字
            ArtTextDesign.add_art_text(
                slide, left=2.5, top=3, width=8, height=2,
                text=f"背景色: RGB{color}",
                shape_type="rounded_rect",
                shape_color=(240, 240, 240),
                text_color=(50, 50, 50)
            )
        
        # 创建带图片背景的幻灯片
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        
        # 设置背景图片，使用更适合的透明度和暗化效果
        BackgroundDesign.set_background_image(
            slide, 
            "assets/qh.jpg",
            transparency=0.35,  # 更轻微的透明度
            overlay_color=(30, 30, 60)  # 深蓝色调的暗化效果
        )
        
        # 添加带背景框的标题，确保可见性
        TextBoxDesign.add_textbox(
            slide, left=1, top=1, width=10, height=1.5,
            text="图片背景测试", level=TextLevel.TITLE,
            color=(255, 255, 255)  # 白色文字
        )
        
        # 添加艺术字（带半透明背景，确保在图片背景上可见）
        ArtTextDesign.add_art_text(
            slide, left=2, top=3, width=9, height=2,
            text="清华校园作为背景图片",
            shape_type="cloud",
            shape_color=(230, 230, 230),
            border_color=(100, 100, 100),
            border_width=2.0,
            text_color=(0, 0, 150)
        )
        
        # 创建带图片背景的幻灯片 - 添加更自然的蓝色滤镜
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        BackgroundDesign.set_background_image(
            slide, 
            "assets/qh.jpg", 
            transparency=0.7,  # 较高的透明度使滤镜效果更细腻
            overlay_color=(180, 200, 255)  # 更柔和的蓝色
        )
        TextBoxDesign.add_textbox(
            slide, left=1, top=1, width=10, height=1.5,
            text="带蓝色滤镜的背景", level=TextLevel.TITLE,
            color=(0, 0, 100)  # 深蓝色文字
        )
        
        designer.save()
        print(f"背景变体测试文件已生成: {output_path}")
    
    def test_image_with_text_overlay(self):
        """测试图片与文字叠加效果"""
        output_path = os.path.join(self.output_dir, "image_with_text_overlay.pptx")
        designer = PPTDesigner(output_path=output_path)
        
        # 添加带背景图片的幻灯片
        slide = designer.prs.slides.add_slide(designer.prs.slide_layouts[6])
        
        # 设置背景图片
        BackgroundDesign.set_background_image(slide, "assets/qh.jpg")
        
        # 添加半透明文本框
        # 注意：此处需要通过形状添加实现半透明效果
        shape = slide.shapes.add_shape(
            1, Inches(2), Inches(2.5), Inches(9), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # 白色
        shape.fill.transparency = 0.5  # 50%透明度
        
        # 添加文本
        shape.text = "图片与文字叠加效果测试"
        text_frame = shape.text_frame
        text_frame.paragraphs[0].alignment = 1  # 居中对齐
        # 修复：使用Pt对象而非直接使用整数
        text_frame.paragraphs[0].font.size = Pt(32)  # 字号
        text_frame.paragraphs[0].font.name = "微软雅黑"
        text_frame.paragraphs[0].font.bold = True
        
        designer.save()
        print(f"图片与文字叠加效果测试文件已生成: {output_path}")

if __name__ == "__main__":
    unittest.main()
