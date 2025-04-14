from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd
from config import STYLES, LAYOUTS, COLORS
from PIL import Image, ImageDraw
from io import BytesIO
import os
import logging
import math

class ProfessionalPPTBuilder:
    def __init__(self, template_path=None, theme='default'):
        """
        专业级PPT构建器
        
        参数:
            template_path (str): 自定义模板路径
            theme (str): 主题名称 (default/blue/creative)
        """
        self.logger = self._setup_logger()
        self.theme = theme
        self.template_path = template_path or self._get_default_template()
        self._ensure_template_exists(self.template_path)
        self.prs = Presentation(self.template_path)
        self.current_section = None
        self.sections = []
        self.logger.info(f"专业PPT构建器初始化完成，使用主题: {theme}")

    def _setup_logger(self):
        """配置专业日志记录"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        return logging.getLogger('ProfessionalPPTBuilder')

    def _get_default_template(self):
        """获取当前主题的默认模板路径"""
        return os.path.join('templates', f'{self.theme}_template.pptx')

    def _ensure_template_exists(self, path):
        """创建专业模板（如果不存在）"""
        if not os.path.exists(path):
            self.logger.warning(f"专业模板不存在，正在创建: {path}")
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            prs = Presentation()
            self._setup_master_slides(prs)
            self._create_default_layouts(prs)
            prs.save(path)
            self.logger.info(f"专业模板创建完成: {path}")

    def _setup_master_slides(self, prs):
        """设置母版样式"""
        # 设置背景和主题色
        background = prs.slide_master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['background'])
        
        # 设置标题样式
        title_placeholder = prs.slide_master.placeholders[0]  # 通常第一个占位符是标题
        title_placeholder.text_frame.text = "[标题样式]"
        title_placeholder.text_frame.paragraphs[0].font.name = STYLES['cover_title']['font']
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(STYLES['cover_title']['size'])
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(*STYLES['cover_title']['color'])
        title_placeholder.text_frame.paragraphs[0].font.bold = STYLES['cover_title']['bold']

        # 设置正文样式
        body_placeholder = prs.slide_master.placeholders[1]  # 通常第二个占位符是正文
        body_placeholder.text_frame.text = "[正文样式]"
        body_placeholder.text_frame.paragraphs[0].font.name = STYLES['default']['font']
        body_placeholder.text_frame.paragraphs[0].font.size = Pt(STYLES['default']['size'])
        body_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(*STYLES['default']['color'])

    def _create_default_layouts(self, prs):
        """创建专业布局"""
        # 使用现有的布局
        # 1. 设计封面布局 - 使用第一个布局作为封面
        if len(prs.slide_layouts) > 0:
            cover_layout = prs.slide_layouts[0]
            self._design_basic_layout(cover_layout, title="[主标题]", name="专业封面")
        
        # 2. 设计目录布局 - 使用第二个布局作为目录页
        if len(prs.slide_layouts) > 1:
            toc_layout = prs.slide_layouts[1]
            self._design_basic_layout(toc_layout, title="目录", name="目录页")
        
        # 3. 设计章节过渡页布局 - 使用第三个布局作为章节页
        if len(prs.slide_layouts) > 2:
            section_layout = prs.slide_layouts[2]
            self._design_basic_layout(section_layout, title="[章节标题]", name="章节页")
        
        # 4. 设计图表布局 - 使用第四个布局作为图表页
        if len(prs.slide_layouts) > 3:
            chart_layout = prs.slide_layouts[3]
            self._design_basic_layout(chart_layout, title="[图表标题]", name="图表页")
        
        # 注意：python-pptx 不支持直接通过代码添加新的幻灯片布局（SlideLayouts），
        # 我们只能使用已有的布局进行自定义。
        
        # 如果没有足够的布局，提供警告
        if len(prs.slide_layouts) < 4:
            self.logger.warning(f"模板中只有 {len(prs.slide_layouts)} 个布局，可能无法满足所有需求。")

    def _design_basic_layout(self, layout, title="", name=""):
        """设计基础布局，为所有布局提供统一的基础配置"""
        layout.name = name
        
        # 清理现有占位符
        for placeholder in list(layout.placeholders):
            try:
                # 保留标题和正文占位符（通常为前两个）
                if placeholder.placeholder_format.idx in [0, 1]:
                    continue
                placeholder.element.getparent().remove(placeholder.element)
            except:
                pass
        
        # 确保有标题占位符
        if len(layout.placeholders) == 0:
            self.logger.warning(f"无法在布局中找到标题占位符: {name}")
        else:
            # 使用第一个占位符作为标题
            try:
                title_placeholder = layout.placeholders[0]
                title_placeholder.text = title
                
                # 应用标题样式
                for paragraph in title_placeholder.text_frame.paragraphs:
                    paragraph.font.name = STYLES['content_title']['font']
                    paragraph.font.size = Pt(STYLES['content_title']['size'])
                    paragraph.font.color.rgb = RGBColor(*STYLES['content_title']['color'])
                    paragraph.font.bold = STYLES['content_title']['bold']
            except Exception as e:
                self.logger.error(f"设置标题失败: {e}")
                
        # 注意：无法直接在布局上添加装饰线，这将在幻灯片实例化时添加

    def _design_cover_layout(self, layout):
        """设计专业封面布局"""
        try:
            # 检查主标题占位符是否存在
            title = layout.placeholders[10]  # 主标题
            title.text = "[主标题]"
            title.left, title.top = Emu(1 * 360000), Emu(2 * 360000)  # 1英寸, 2英寸
            title.width, title.height = Emu(8 * 360000), Emu(1.5 * 360000)

            # 检查副标题占位符是否存在
            subtitle = layout.placeholders[11]  # 副标题
            subtitle.text = "[副标题]"
            subtitle.left, subtitle.top = Emu(1 * 360000), Emu(3.5 * 360000)

            # 添加装饰元素
            line = layout.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Emu(1 * 360000), Emu(3.2 * 360000), 
                Emu(3 * 360000), Emu(0.05 * 360000)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['accent'])
            line.line.fill.background()
        except IndexError as e:
            self.logger.error("封面布局缺少必要的占位符，请检查模板文件是否包含至少 11 个占位符。")
            raise e

    def _design_toc_layout(self, layout):
        """设计目录页布局"""
        title = layout.placeholders[10]
        title.text = "目录"
        title.left, title.top = Emu(0.5 * 360000), Emu(0.5 * 360000)
        
        # 添加目录占位框
        toc_box = layout.placeholders.add_placeholder(
            11, 
            Emu(1 * 360000), Emu(1.5 * 360000),
            Emu(8 * 360000), Emu(4 * 360000)
        )
        toc_box.text_frame.word_wrap = True
        toc_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    def _design_section_layout(self, layout):
        """设计章节过渡页"""
        title = layout.placeholders[10]
        title.text = "[章节标题]"
        title.left, title.top = Emu(1 * 360000), Emu(1 * 360000)
        
        # 添加章节进度指示器
        for i in range(5):
            dot = layout.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu((1+i*1.5)*360000), Emu(3 * 360000),
                Emu(0.8 * 360000), Emu(0.8 * 360000)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['secondary'])
            if i == 0:
                dot.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['primary'])

    def _design_content_layouts(self, prs):
        """设计多种内容版式"""
        # 版式1: 标题+正文
        layout1 = prs.slide_layouts.add_slide_layout(prs.slide_master.slide_layouts[1])
        layout1.name = "内容-正文"
        self._design_content_layout(layout1, with_image=False)
        
        # 版式2: 标题+正文+图片
        layout2 = prs.slide_layouts.add_slide_layout(prs.slide_master.slide_layouts[1])
        layout2.name = "内容-图文"
        self._design_content_layout(layout2, with_image=True)

    def _design_content_layout(self, layout, with_image):
        """设计内容页基础布局"""
        title = layout.placeholders[10]
        title.text = "[内容标题]"
        title.left, title.top = Emu(0.5 * 360000), Emu(0.5 * 360000)
        
        # 正文框
        content = layout.placeholders.add_placeholder(
            11, 
            Emu(0.5 * 360000), Emu(1.5 * 360000),
            Emu(5 * 360000) if with_image else Emu(8 * 360000), 
            Emu(4 * 360000)
        )
        content.text_frame.word_wrap = True
        
        # 图片占位区
        if with_image:
            img_placeholder = layout.placeholders.add_placeholder(
                12,
                Emu(6 * 360000), Emu(1.5 * 360000),
                Emu(3 * 360000), Emu(4 * 360000)
            )
            img_placeholder.name = "IMAGE_PLACEHOLDER"

    def _design_chart_layout(self, layout):
        """设计图表页布局"""
        title = layout.placeholders[10]
        title.text = "[图表标题]"
        title.left, title.top = Emu(0.5 * 360000), Emu(0.5 * 360000)
        
        chart_area = layout.placeholders.add_placeholder(
            11,
            Emu(1 * 360000), Emu(1.5 * 360000),
            Emu(8 * 360000), Emu(4 * 360000)
        )
        chart_area.name = "CHART_AREA"

    def add_cover(self, title, subtitle, logo_path=None, date=None, author=None):
        """添加专业封面页"""
        try:
            slide = self.prs.slides.add_slide(self._get_layout("专业封面"))
            
            # 设置主标题
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            self._apply_style(title_placeholder, "cover_title")
            
            # 设置副标题
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
            self._apply_style(subtitle_placeholder, "cover_subtitle")
            
            # 添加LOGO (仅当文件确实存在时)
            if logo_path:
                # 使用用户提供的路径
                if os.path.exists(logo_path):
                    try:
                        slide.shapes.add_picture(
                            logo_path, 
                            Emu(7.5 * 360000), Emu(0.5 * 360000),
                            height=Emu(1 * 360000)
                        )
                    except Exception as e:
                        self.logger.warning(f"无法添加LOGO图像 {logo_path}: {e}")
            else:
                # 使用绝对路径访问默认logo
                script_dir = os.path.dirname(os.path.abspath(__file__))
                default_logo = os.path.join(script_dir, 'assets', 'real_logo.png')
                if os.path.exists(default_logo):
                    try:
                        slide.shapes.add_picture(
                            default_logo, 
                            Emu(7.5 * 360000), Emu(0.5 * 360000),
                            height=Emu(1 * 360000)
                        )
                    except Exception as e:
                        self.logger.warning(f"无法添加默认LOGO图像: {e}")
                        # 尝试备用logo
                        backup_logo = os.path.join(script_dir, 'assets', 'placeholder_logo.png')
                        if os.path.exists(backup_logo):
                            try:
                                slide.shapes.add_picture(
                                    backup_logo,
                                    Emu(7.5 * 360000), Emu(0.5 * 360000),
                                    height=Emu(1 * 360000)
                                )
                            except Exception as e:
                                self.logger.warning(f"无法添加备用LOGO图像: {e}")
            
            # 添加页脚信息
            if date or author:
                footer_text = []
                if author:
                    footer_text.append(f"制作: {author}")
                if date:
                    footer_text.append(f"日期: {date}")
                
                footer = slide.shapes.add_textbox(
                    Emu(0.5 * 360000), Emu(6.8 * 360000),
                    Emu(9 * 360000), Emu(0.5 * 360000)
                )
                footer.text_frame.text = " | ".join(footer_text)
                self._apply_style(footer, "footer")
            
            self.logger.info("专业封面页添加完成")
            return slide
        except Exception as e:
            self.logger.error(f"添加封面页失败: {e}")
            raise

    def add_toc(self, sections, style="default"):
        """添加目录页"""
        try:
            slide = self.prs.slides.add_slide(self._get_layout("目录页"))
            title = slide.shapes.title
            title.text = "目录"
            self._apply_style(title, "section_title")
            
            toc_box = slide.placeholders[1]
            toc_text = "\n".join([f"{i+1}. {sec['title']}" for i, sec in enumerate(sections)])
            toc_box.text_frame.text = toc_text
            
            # 自动调整字号避免溢出
            self._auto_adjust_text_size(toc_box.text_frame, max_lines=len(sections)*2)
            
            self.sections = sections
            self.logger.info("目录页添加完成")
            return slide
        except Exception as e:
            self.logger.error(f"添加目录页失败: {e}")
            raise

    def add_section(self, title, index=None):
        """添加章节过渡页"""
        try:
            slide = self.prs.slides.add_slide(self._get_layout("章节页"))
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_style(title_shape, "section_title")
            
            # 更新当前章节
            self.current_section = title
            
            # 更新进度指示器
            if index is not None and index < len(self.sections):
                for i, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE.OVAL:
                        if i-2 == index:  # 假设前两个形状是标题和装饰
                            shape.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['primary'])
                        else:
                            shape.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['secondary'])
            
            self.logger.info(f"章节页添加完成: {title}")
            return slide
        except Exception as e:
            self.logger.error(f"添加章节页失败: {e}")
            raise

    def add_content_slide(self, title, content, layout_type="正文", image_path=None):
        """添加专业内容页"""
        try:
            layout_name = f"内容-{layout_type}"
            slide = self.prs.slides.add_slide(self._get_layout(layout_name))
            
            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_style(title_shape, "content_title")
            
            # 设置内容
            content_box = slide.placeholders[1]
            content_box.text_frame.text = content
            self._apply_style(content_box, "content_text")
            
            # 自动调整文本大小
            self._auto_adjust_text_size(content_box.text_frame)
            
            # 添加图片
            if image_path and os.path.exists(image_path) and "图文" in layout_type:
                img_placeholder = next((s for s in slide.shapes if s.name == "IMAGE_PLACEHOLDER"), None)
                if img_placeholder:
                    img_placeholder.element.getparent().remove(img_placeholder.element)
                    slide.shapes.add_picture(
                        image_path,
                        img_placeholder.left, img_placeholder.top,
                        width=img_placeholder.width, height=img_placeholder.height
                    )
            
            self.logger.info(f"内容页添加完成: {title}")
            return slide
        except Exception as e:
            self.logger.error(f"添加内容页失败: {e}")
            raise

    def add_chart_slide(self, title, chart_data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
        """添加专业图表页"""
        try:
            slide = self.prs.slides.add_slide(self._get_layout("图表页"))
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_style(title_shape, "chart_title")
            
            # 准备图表数据
            chart_data_obj = ChartData()
            chart_data_obj.categories = chart_data["categories"]
            for series in chart_data["series"]:
                chart_data_obj.add_series(series["name"], series["values"])
            
            # 添加图表
            chart_area = next((s for s in slide.shapes if s.name == "CHART_AREA"), None)
            if chart_area:
                chart_area.element.getparent().remove(chart_area.element)
                chart = slide.shapes.add_chart(
                    chart_type, 
                    chart_area.left, chart_area.top,
                    chart_area.width, chart_area.height,
                    chart_data_obj
                ).chart
                
                # 专业图表样式
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                chart.plots[0].has_data_labels = True
                
                # 应用主题颜色
                for i, series in enumerate(chart.series):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(*COLORS[self.theme]['chart_colors'][i % 3])
            
            self.logger.info(f"图表页添加完成: {title}")
            return slide
        except Exception as e:
            self.logger.error(f"添加图表页失败: {e}")
            raise

    def _get_layout(self, layout_name):
        """获取指定名称的布局"""
        for layout in self.prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return self.prs.slide_layouts[1]  # 默认返回标题布局

    def _apply_style(self, shape, style_type):
        """应用专业样式"""
        style = STYLES.get(style_type, STYLES["default"])
        
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                if "font" in style:
                    paragraph.font.name = style["font"]
                if "size" in style:
                    paragraph.font.size = Pt(style["size"])
                if "color" in style:
                    paragraph.font.color.rgb = RGBColor(*style["color"])
                if "bold" in style:
                    paragraph.font.bold = style["bold"]
                if "italic" in style:
                    paragraph.font.italic = style["italic"]
                if "underline" in style:
                    paragraph.font.underline = style["underline"]
                if "align" in style:
                    paragraph.alignment = getattr(PP_PARAGRAPH_ALIGNMENT, style["align"].upper())
                
                # 行距设置
                if "line_spacing" in style:
                    paragraph.line_spacing = style["line_spacing"]
                
                # 段前段后间距
                if "space_before" in style:
                    paragraph.space_before = Pt(style["space_before"])
                if "space_after" in style:
                    paragraph.space_after = Pt(style["space_after"])

    def _auto_adjust_text_size(self, text_frame, max_lines=10, min_size=12):
        """自动调整文本大小避免溢出"""
        original_size = text_frame.paragraphs[0].font.size.pt if text_frame.paragraphs else 14
        
        # 计算需要的行数
        line_count = sum([len(p.text.split('\n')) for p in text_frame.paragraphs])
        
        # 如果超出最大行数，减小字号
        if line_count > max_lines:
            ratio = max_lines / line_count
            new_size = max(min_size, math.floor(original_size * ratio))
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(new_size)

    def save(self, output_path=None):
        """保存专业PPT"""
        try:
            output_path = output_path or os.path.join('output', f'professional_{self.theme}_presentation.pptx')
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)
            self.logger.info(f"专业PPT保存成功: {output_path}")
            return output_path
        except Exception as e:
            self.logger.error(f"保存PPT失败: {e}")
            raise

    def build_demo_presentation(self):
        """构建完整演示示例"""
        try:
            self.logger.info("开始构建专业演示文稿示例")
            
            # 1. 添加封面
            self.add_cover(
                title="2023年度业务报告",
                subtitle="市场分析与未来展望",
                author="数据分析部",
                date="2023年12月"
            )
            
            # 2. 添加目录
            sections = [
                {"title": "市场概况", "slides": 3},
                {"title": "销售分析", "slides": 2},
                {"title": "未来计划", "slides": 2}
            ]
            self.add_toc(sections)
            
            # 3. 添加各章节
            for i, section in enumerate(sections):
                self.add_section(section["title"], index=i)
                
                # 添加内容页
                if section["title"] == "市场概况":
                    self.add_content_slide(
                        title="市场趋势",
                        content="2023年市场呈现以下特点：\n\n• 整体规模增长15%\n• 新兴领域增速达40%\n• 传统行业面临转型\n• 消费者偏好变化显著",
                        layout_type="图文",
                        image_path=os.path.join('assets', 'market_trend.png')
                    )
                    
                    # 添加图表
                    market_data = {
                        "categories": ["Q1", "Q2", "Q3", "Q4"],
                        "series": [
                            {"name": "国内市场", "values": [120, 135, 145, 165]},
                            {"name": "国际市场", "values": [80, 90, 110, 125]}
                        ]
                    }
                    self.add_chart_slide(
                        title="季度市场分布",
                        chart_data=market_data
                    )
                
                elif section["title"] == "销售分析":
                    # 添加更多专业内容...
                    pass
            
            # 保存演示文稿
            output_file = self.save()
            self.logger.info(f"专业演示文稿构建完成: {output_file}")
            return output_file
        except Exception as e:
            self.logger.error(f"构建演示示例失败: {e}")
            raise


if __name__ == "__main__":
    # 示例用法
    builder = ProfessionalPPTBuilder(theme='blue')
    
    # 构建完整示例演示
    output_file = builder.build_demo_presentation()
    print(f"专业演示文稿已生成: {output_file}")