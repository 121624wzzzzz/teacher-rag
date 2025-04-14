from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData  # 新增导入
from pptx.enum.chart import XL_CHART_TYPE  # 确保已导入

def create_ppt_template(output_path="professional_template.pptx"):
    # 初始化演示文稿
    prs = Presentation()
    
    # ===== 1. 定义主题配色和字体 =====
    prs.slide_width = Inches(13.33)  # 16:9比例
    prs.slide_height = Inches(7.5)
    
    # 主色：深蓝 (#2A5C82)，辅助色：浅蓝 (#4EACCD)，强调色：红 (#FF6B6B)
    colors = {
        "primary": RGBColor(42, 92, 130),
        "secondary": RGBColor(78, 172, 205),
        "accent": RGBColor(255, 107, 107)
    }
    
    # ===== 2. 设计母版版式 =====
    slide_master = prs.slide_master
    
    # 设置母版背景（渐变）
    background = slide_master.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = colors["primary"]
    fill.gradient_stops[1].color.rgb = RGBColor(247, 249, 252)  # 浅灰
    
    # 设置母版标题和正文字体
    title_style = slide_master.slide_layouts[0].placeholders[0].text_frame.paragraphs[0].font
    title_style.name = "微软雅黑"
    title_style.size = Pt(44)
    title_style.bold = True
    
    body_style = slide_master.slide_layouts[1].placeholders[1].text_frame.paragraphs[0].font
    body_style.name = "思源黑体"
    body_style.size = Pt(24)
    
    # ===== 3. 创建核心页面 =====
    
    # ---- 封面页 ----
    slide_cover = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    title = slide_cover.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title.text = "演示文稿标题"
    title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = slide_cover.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle.text = "副标题/日期/作者"
    subtitle.text_frame.paragraphs[0].font.color.rgb = colors["secondary"]
    
    # ---- 目录页 ----
    slide_toc = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    title = slide_toc.shapes.title
    title.text = "目录"
    
    content = slide_toc.placeholders[1].text_frame
    for item in ["市场分析", "产品介绍", "数据报告", "行动计划"]:
        p = content.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(28)
    
    # ---- 图文内容页（左文右图） ----
    slide_content = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    
    # 左侧文字
    left_text = slide_content.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(5))
    tf = left_text.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "核心观点标题"
    p.font.bold = True
    p.font.size = Pt(32)
    
    for i in range(3):
        p = tf.add_paragraph()
        p.text = f"这里是第{i+1}个要点内容，支持自动换行"
        p.level = 1
        p.font.size = Pt(24)
    
    # 右侧图片占位符
    right_pic = slide_content.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1), Inches(6), Inches(4.5)
    )
    right_pic.text = "[图片占位符]"
    right_pic.fill.solid()
    right_pic.fill.fore_color.rgb = RGBColor(230, 230, 230)  # 浅灰占位
    
    # ---- 数据图表页 ----
    slide_chart = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide_chart.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.8))
    title.text = "年度数据报告"
    
    # 修正：使用 ChartData 对象
    chart_data = ChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]  # X轴标签
    chart_data.add_series('销售额', (12.5, 18.3, 15.7, 21.4))  # Y轴数据
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
    chart = slide_chart.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # 设置图表颜色
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = colors["accent"]
    
    # ===== 4. 保存文件 =====
    prs.save(output_path)
    print(f"模板已生成: {output_path}")

if __name__ == "__main__":
    create_ppt_template()